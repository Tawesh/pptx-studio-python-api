# -*- coding: utf-8 -*-
"""渲染引擎：演示文稿创建、页面编排、组件渲染、导出。"""
import os
import uuid
import logging
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from . import utils as U
from . import components as C
from .themes import get_theme

logger = logging.getLogger("pptx_studio")

DEFAULT_FILENAME_PREFIX = "pptx_studio"


def _gen_filename(prefix, ext="pptx"):
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    uid = uuid.uuid4().hex[:8]
    return f"{prefix}_{ts}_{uid}.{ext}"


def _parse_size(size_spec):
    """解析画布尺寸，返回 (width_emu, height_emu)。"""
    if isinstance(size_spec, dict):
        return Inches(float(size_spec.get("w", 13.333))), Inches(float(size_spec.get("h", 7.5)))
    s = str(size_spec or "16:9").lower()
    if s in ("4:3", "4_3"):
        return Inches(10), Inches(7.5)
    return Inches(13.333), Inches(7.5)


def render(data: dict, output_dir: str) -> str:
    """渲染完整 PPTX，返回文件名。

    data 结构：
    {
      "theme": "deep_space" | dict,     # 必填（默认 deep_space）
      "size": "16:9" | "4:3" | {"w":..,"h":..},
      "filename_prefix": "xxx",
      "pages": [
        {
          "type": "cover|section|content|blank",
          "title": "...",
          "subtitle": "...",
          "background": {...},
          "decorations": [...],
          "components": [...],
          "page_number": true,
          "footer": "..."
        }
      ]
    }
    """
    theme = get_theme(data.get("theme", "deep_space"))
    prs = Presentation()
    prs.slide_width, prs.slide_height = _parse_size(data.get("size", "16:9"))
    sw, sh = prs.slide_width, prs.slide_height

    ctx = {
        "prs": prs,
        "slide_width": sw,
        "slide_height": sh,
        "fonts": theme.get("fonts", {"heading": "微软雅黑", "body": "微软雅黑"}),
        "style": theme.get("style", "dark"),
    }

    pages = data.get("pages", [])
    if not pages and data.get("title"):
        # 兼容旧式：单页 cover
        pages = [{"type": "cover", "title": data.get("title"),
                  "subtitle": data.get("subtitle")}]

    page_index = 0
    for page in pages:
        page_index += 1
        _render_page(prs, page, theme, ctx, page_index, len(pages))

    os.makedirs(output_dir, exist_ok=True)
    prefix = _safe_prefix(data.get("filename_prefix", DEFAULT_FILENAME_PREFIX))
    filename = _gen_filename(prefix)
    filepath = os.path.join(output_dir, filename)
    prs.save(filepath)
    return filename


def _safe_prefix(prefix):
    return "".join(c for c in str(prefix) if c.isalnum() or c in "-_").strip() or DEFAULT_FILENAME_PREFIX


def _render_page(prs, page, theme, ctx, page_index, total_pages):
    """渲染单个页面。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    # 1) 背景
    U.apply_background(slide, page.get("background"), theme, prs)

    # 2) 装饰
    U.apply_decorations(slide, page.get("decorations"), theme, prs)

    ptype = str(page.get("type", "content")).lower()

    if ptype == "cover":
        _render_cover(slide, page, theme, ctx)
    elif ptype == "section":
        _render_section(slide, page, theme, ctx)
    elif ptype == "content":
        _render_content_header(slide, page, theme, ctx)

    # 3) 自由组件（内容页/空白页/封面附加组件）
    comps = page.get("components", [])
    if ptype == "cover" and comps:
        # 封面组件：相对底部区域渲染
        for spec in comps:
            C.render_component(slide, spec, theme, ctx)
    elif ptype in ("content", "blank"):
        for spec in comps:
            C.render_component(slide, spec, theme, ctx)

    # 4) 页脚 / 页码
    footer = page.get("footer")
    if footer:
        C.render_footer(slide, {"type": "footer", "text": footer,
                                "page": page.get("page_number", False),
                                "page_number": page_index}, theme, ctx)
    elif page.get("page_number", False):
        C.render_page_number(slide, {"type": "page_number", "number": page_index}, theme, ctx)


def _render_cover(slide, page, theme, ctx):
    """封面：居中大标题 + 副标题 + 底部标签。"""
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    sw = ctx["slide_width"]
    title = page.get("title", "")
    subtitle = page.get("subtitle", "")
    title_color = U.color_hex(page.get("title_color", "text"), theme)
    accent = U.color_hex(page.get("accent", "primary"), theme)

    # 顶部小徽标（可选）
    if page.get("kicker"):
        kicker_w = Inches(2.2)
        kicker_x = Emu(int((sw - kicker_w) / 2))
        badge = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, kicker_x, Inches(2.0), kicker_w, Inches(0.5))
        U.set_shape_fill(badge, accent, 18)
        badge.line.fill.background()
        U.no_shadow(badge)
        btf = badge.text_frame
        btf.word_wrap = False
        btf.vertical_anchor = MSO_ANCHOR.MIDDLE
        btf.margin_left = btf.margin_right = 0
        btf.margin_top = btf.margin_bottom = 0
        bp = btf.paragraphs[0]
        bp.alignment = PP_ALIGN.CENTER
        br = bp.add_run()
        U.style_run(br, text=page["kicker"], size=13, bold=True, color=accent,
                    font=ctx["fonts"]["heading"])
        title_y = Inches(2.6)
    else:
        title_y = Inches(2.55)

    # 主标题
    tf = U.textbox(slide, Inches(1.0), title_y, sw - Inches(2.0), Inches(1.5),
                   anchor=MSO_ANCHOR.MIDDLE)
    p = U.add_para(tf, first=True, align=PP_ALIGN.CENTER)
    r = p.add_run()
    U.style_run(r, text=title, size=page.get("title_size", 48), bold=True,
                color=title_color, font=ctx["fonts"]["heading"])

    # 副标题
    if subtitle:
        stf = U.textbox(slide, Inches(2.0), Inches(4.1), sw - Inches(4.0), Inches(0.6),
                        anchor=MSO_ANCHOR.TOP)
        sp = U.add_para(stf, first=True, align=PP_ALIGN.CENTER)
        sr = sp.add_run()
        U.style_run(sr, text=subtitle, size=18, color=accent, font=ctx["fonts"]["body"])

    # 底部装饰线
    line = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Emu(int((sw - Inches(2)) / 2)), Inches(4.9), Inches(2), Inches(0.06))
    U.set_shape_fill(line, accent)
    U.set_round_rect_radius(line, 0.5)
    line.line.fill.background()
    U.no_shadow(line)


def _render_section(slide, page, theme, ctx):
    """章节页：居中大标题 + 序号 + 副标题。"""
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    sw = ctx["slide_width"]
    title = page.get("title", "")
    subtitle = page.get("subtitle", "")
    index = page.get("index")
    accent = U.color_hex(page.get("accent", "primary"), theme)

    y0 = Inches(2.4)
    if index:
        itf = U.textbox(slide, Inches(1.0), y0, sw - Inches(2.0), Inches(1.0),
                        anchor=MSO_ANCHOR.MIDDLE)
        ip = U.add_para(itf, first=True, align=PP_ALIGN.CENTER)
        ir = ip.add_run()
        U.style_run(ir, text=str(index), size=64, bold=True, color=accent,
                    font=ctx["fonts"]["heading"])
        y0 = Inches(3.3)
    tf = U.textbox(slide, Inches(1.0), y0, sw - Inches(2.0), Inches(1.1),
                   anchor=MSO_ANCHOR.MIDDLE)
    p = U.add_para(tf, first=True, align=PP_ALIGN.CENTER)
    r = p.add_run()
    U.style_run(r, text=title, size=page.get("title_size", 40), bold=True,
                color=U.color_hex(page.get("title_color", "text"), theme),
                font=ctx["fonts"]["heading"])
    if subtitle:
        stf = U.textbox(slide, Inches(2.5), y0 + Inches(1.15), sw - Inches(5.0), Inches(0.5))
        sp = U.add_para(stf, first=True, align=PP_ALIGN.CENTER)
        sr = sp.add_run()
        U.style_run(sr, text=subtitle, size=15,
                    color=U.color_hex(page.get("subtitle_color", "text_muted"), theme),
                    font=ctx["fonts"]["body"])


def _render_content_header(slide, page, theme, ctx):
    """内容页顶部标题区：小标题 + 主标题 + 分隔。"""
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    sw = ctx["slide_width"]
    title = page.get("title", "")
    kicker = page.get("kicker")
    accent = U.color_hex(page.get("accent", "primary"), theme)
    title_color = U.color_hex(page.get("title_color", "text"), theme)

    title_x = Inches(0.7)
    if kicker:
        ktf = U.textbox(slide, title_x, Inches(0.35), sw - Inches(2.0), Inches(0.35))
        kp = U.add_para(ktf, first=True)
        kr = kp.add_run()
        U.style_run(kr, text=kicker.upper(), size=12, bold=True, color=accent,
                    font=ctx["fonts"]["heading"])
        title_y = Inches(0.72)
    else:
        title_y = Inches(0.45)
    if title:
        tf = U.textbox(slide, title_x, title_y, sw - Inches(2.0), Inches(0.8),
                       anchor=MSO_ANCHOR.MIDDLE)
        p = U.add_para(tf, first=True)
        r = p.add_run()
        U.style_run(r, text=title, size=page.get("title_size", 30), bold=True,
                    color=title_color, font=ctx["fonts"]["heading"])
    # 标题下方强调条
    bar = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, title_x, Inches(1.35), Inches(0.6), Inches(0.07))
    U.set_shape_fill(bar, accent)
    U.set_round_rect_radius(bar, 0.5)
    bar.line.fill.background()
    U.no_shadow(bar)
