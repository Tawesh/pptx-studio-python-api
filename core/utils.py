# -*- coding: utf-8 -*-
"""底层渲染工具：颜色/字体/渐变/圆角/阴影/尺寸解析/背景/装饰。"""
import math
import os

from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn


# ------------------------------------------------------------------ 颜色
def color_hex(v, theme=None):
    """把组件传入的颜色值解析为 6 位 HEX。
    支持：
      - None            -> 主题 primary
      - "primary"/"accent" 等主题键
      - "#RRGGBB" / "RRGGBB"
      - (r,g,b) 元组
    """
    if v is None:
        return (theme or {}).get("colors", {}).get("primary", "1F4E79")
    if isinstance(v, (tuple, list)):
        return "%02X%02X%02X" % tuple(int(x) for x in v)
    s = str(v).strip().lstrip("#").upper()
    if s in ("PRIMARY", "SECONDARY", "ACCENT", "TEXT", "TEXT_MUTED", "TEXT_INVERT",
             "SUCCESS", "WARNING", "DANGER", "CARD_BG", "CARD_BG2", "CARD_BORDER", "OVERLAY"):
        return (theme or {}).get("colors", {}).get(s.lower(), s)
    if len(s) == 6 and all(c in "0123456789ABCDEF" for c in s):
        return s
    return "1F4E79"


def alpha_hex(hex_color, alpha):
    """返回带透明度的 XML 颜色元素字符串（alpha 0-100）。"""
    return f'<a:srgbClr val="{hex_color}"><a:alpha val="{int(alpha * 1000)}"/></a:srgbClr>'


def hex_to_rgb(hex_color):
    return RGBColor.from_string(hex_color)


# ------------------------------------------------------------------ 字体
def set_run_font(run, name=None):
    """设置 run 字体，含中文 east-asian 字体（用原生 get_or_add_ea 保证元素顺序正确）。"""
    try:
        if name:
            run.font.name = name
        rPr = run._r.get_or_add_rPr()
        ea = rPr.find(qn("a:ea"))
        if ea is None:
            ea = rPr.get_or_add_ea()
        if name:
            ea.set("typeface", name)
    except Exception:
        pass


def style_run(run, text=None, size=None, bold=None, color=None, italic=None,
              font=None, spacing=None):
    """统一设置 run 样式。"""
    if text is not None:
        run.text = text
    if size is not None:
        run.font.size = Pt(size)
    if bold is not None:
        run.font.bold = bold
    if italic is not None:
        run.font.italic = italic
    if color is not None:
        try:
            run.font.color.rgb = RGBColor.from_string(color)
        except Exception:
            pass
    if font is not None:
        set_run_font(run, font)
    if spacing is not None:
        try:
            run.font._rPr.set("spc", str(int(spacing * 100)))
        except Exception:
            pass


def add_para(tf, first=False, align=None, line_spacing=None, space_before=None,
             space_after=None, indent_level=None):
    """在文本框添加/复用段落并设置段落级样式。"""
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    if align is not None:
        p.alignment = align
    if line_spacing is not None:
        p.line_spacing = line_spacing
    if space_before is not None:
        p.space_before = Pt(space_before)
    if space_after is not None:
        p.space_after = Pt(space_after)
    if indent_level is not None:
        p.level = indent_level
    return p


# ------------------------------------------------------------------ 尺寸
def parse_len(v, container_emu):
    """解析长度：数字(英寸)、'50%'(百分比)、'fill'(填满容器)、'auto'。"""
    if v is None:
        return container_emu
    if isinstance(v, (int, float)):
        return Inches(float(v))
    s = str(v).strip().lower()
    if s == "fill":
        return Emu(int(container_emu))
    if s.endswith("%"):
        try:
            return Emu(int(int(s[:-1]) / 100.0 * container_emu))
        except Exception:
            return container_emu
    if s.endswith("cm"):
        try:
            return Emu(int(float(s[:-2]) * 360000))
        except Exception:
            return container_emu
    try:
        return Inches(float(s))
    except Exception:
        return container_emu


# ------------------------------------------------------------------ 形状
def set_shape_fill(shape, color=None, alpha=None):
    """纯色填充（可选透明度）。"""
    if color is None:
        return
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(color)
    if alpha is not None and 0 <= alpha < 100:
        _set_fill_alpha(shape, alpha)


def _set_fill_alpha(shape, alpha):
    try:
        spPr = shape._element.spPr
        solidFill = spPr.find(qn("a:solidFill"))
        if solidFill is not None:
            clr = solidFill.find(qn("a:srgbClr"))
            if clr is not None:
                a = clr.makeelement(qn("a:alpha"), {"val": str(int(alpha * 1000))})
                clr.append(a)
    except Exception:
        pass


def set_shape_line(shape, color=None, width_pt=None, alpha=None):
    """设置描边。"""
    line = shape.line
    if color is None:
        line.fill.background()
    else:
        line.color.rgb = RGBColor.from_string(color)
        if width_pt is not None:
            line.width = Pt(width_pt)
        if alpha is not None and 0 <= alpha < 100:
            _set_line_alpha(shape, alpha)


def _set_line_alpha(shape, alpha):
    try:
        spPr = shape._element.spPr
        ln = spPr.find(qn("a:ln"))
        if ln is not None:
            solidFill = ln.find(qn("a:solidFill"))
            if solidFill is not None:
                clr = solidFill.find(qn("a:srgbClr"))
                if clr is not None:
                    a = clr.makeelement(qn("a:alpha"), {"val": str(int(alpha * 1000))})
                    clr.append(a)
    except Exception:
        pass


def set_round_rect_radius(shape, radius=0.15):
    """设置圆角矩形圆角半径（0~0.5，0.5 为胶囊）。"""
    try:
        spPr = shape._element.spPr
        geom = spPr.find(qn("a:prstGeom"))
        if geom is None:
            return
        avLst = geom.find(qn("a:avLst"))
        if avLst is None:
            avLst = geom.makeelement(qn("a:avLst"), {})
            geom.append(avLst)
        gd = geom.makeelement(qn("a:gd"),
                              {"name": "adj", "fmla": "val %d" % int(max(0, min(0.5, radius)) * 100000)})
        avLst.append(gd)
    except Exception:
        pass


def set_rounded_picture(shape, radius=0.08):
    """给图片设置圆角裁剪。"""
    try:
        spPr = shape._element.spPr
        geom = spPr.makeelement(qn("a:prstGeom"), {"prst": "roundRect"})
        avLst = spPr.makeelement(qn("a:avLst"), {})
        gd = spPr.makeelement(qn("a:gd"),
                              {"name": "adj", "fmla": "val %d" % int(max(0, min(0.5, radius)) * 100000)})
        avLst.append(gd)
        geom.append(avLst)
        old = spPr.find(qn("a:prstGeom"))
        if old is not None:
            spPr.remove(old)
        # 插入到 xfrm 之后
        xfrm = spPr.find(qn("a:xfrm"))
        if xfrm is not None:
            xfrm.addnext(geom)
        else:
            spPr.insert(0, geom)
    except Exception:
        pass


def add_shadow(shape, blur=0.09, dist=0.04, direction=90, color="000000", alpha=35):
    """给形状添加外阴影。"""
    try:
        spPr = shape._element.spPr
        eff = spPr.find(qn("a:effectLst"))
        if eff is None:
            eff = spPr.makeelement(qn("a:effectLst"), {})
            spPr.append(eff)
        shdw = spPr.makeelement(qn("a:outerShdw"), {
            "blurRad": str(int(blur * 914400)),
            "dist": str(int(dist * 914400)),
            "dir": str(int(direction * 60000)),
            "rotWithShape": "0",
        })
        clr = spPr.makeelement(qn("a:srgbClr"), {"val": color})
        a = spPr.makeelement(qn("a:alpha"), {"val": str(int(alpha * 1000))})
        clr.append(a)
        shdw.append(clr)
        eff.append(shdw)
    except Exception:
        pass


def no_shadow(shape):
    try:
        shape.shadow.inherit = False
    except Exception:
        pass


# ------------------------------------------------------------------ 渐变
def make_gradient_fill(shape, colors, angle=90):
    """给形状应用多色线性渐变（colors 为 HEX 列表）。"""
    try:
        fill = shape.fill
        fill.gradient()
        stops = fill.gradient_stops
        n = len(colors)
        for i, c in enumerate(colors):
            if i < len(stops):
                stops[i].color.rgb = RGBColor.from_string(c)
                stops[i].position = i / (n - 1) if n > 1 else 0.0
        spPr = shape._element.spPr
        gradFill = spPr.find(qn("a:gradFill"))
        if gradFill is not None:
            lin = gradFill.find(qn("a:lin"))
            if lin is None:
                lin = gradFill.makeelement(qn("a:lin"), {})
                gradFill.append(lin)
            lin.set("ang", str(int(angle * 60000)))
            lin.set("scaled", "1")
    except Exception:
        pass


def set_slide_gradient(slide, colors, angle=90):
    """幻灯片背景渐变。"""
    try:
        fill = slide.background.fill
        fill.gradient()
        stops = fill.gradient_stops
        n = len(colors)
        for i, c in enumerate(colors):
            if i < len(stops):
                stops[i].color.rgb = RGBColor.from_string(c)
                stops[i].position = i / (n - 1) if n > 1 else 0.0
        bg = slide.background._element
        bgPr = bg.find(qn("p:bgPr"))
        if bgPr is not None:
            gradFill = bgPr.find(qn("a:gradFill"))
            if gradFill is not None:
                lin = gradFill.find(qn("a:lin"))
                if lin is None:
                    lin = gradFill.makeelement(qn("a:lin"), {})
                    gradFill.append(lin)
                lin.set("ang", str(int(angle * 60000)))
                lin.set("scaled", "1")
    except Exception:
        pass


def set_slide_solid(slide, color):
    try:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor.from_string(color)
    except Exception:
        pass


def apply_background(slide, bg_spec, theme, prs):
    """应用页面背景。bg_spec 可省略（用主题默认）。"""
    if bg_spec is None:
        bg_spec = theme.get("background", {"type": "solid", "color": "FFFFFF"})
    btype = str(bg_spec.get("type", "gradient")).lower()
    colors = bg_spec.get("colors")
    if colors:
        colors = [color_hex(c, theme) for c in colors]
    if btype == "gradient" and colors:
        set_slide_gradient(slide, colors, float(bg_spec.get("angle", 90)))
    elif btype == "solid":
        set_slide_solid(slide, color_hex(bg_spec.get("color", "FFFFFF"), theme))
    elif btype == "image":
        path = bg_spec.get("path")
        if path and os.path.isfile(path):
            try:
                sw = prs.slide_width
                sh = prs.slide_height
                pic = slide.shapes.add_picture(path, 0, 0, width=sw, height=sh)
                overlay = bg_spec.get("overlay", 0.0)
                if overlay:
                    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, sw, sh)
                    rect.fill.solid()
                    rect.fill.fore_color.rgb = RGBColor.from_string(color_hex(theme["colors"]["overlay"], theme))
                    _set_fill_alpha(rect, overlay * 100)
                    rect.line.fill.background()
                    no_shadow(rect)
            except Exception:
                pass


# ------------------------------------------------------------------ 装饰
def add_deco_bar(slide, prs, color, top=True, height=0.08, alpha=100):
    """顶部/底部装饰条。"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, Emu(0 if top else int(prs.slide_height - Inches(height))),
        prs.slide_width, Inches(height),
    )
    set_shape_fill(shape, color, alpha)
    shape.line.fill.background()
    no_shadow(shape)
    return shape


def add_deco_circle(slide, prs, color, x_frac, y_frac, size, alpha=100):
    """装饰圆。"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Emu(int(prs.slide_width * x_frac)), Emu(int(prs.slide_height * y_frac)),
        Inches(size), Inches(size),
    )
    set_shape_fill(shape, color, alpha)
    shape.line.fill.background()
    no_shadow(shape)
    return shape


def add_dot_grid(slide, prs, color, alpha=18, density=0.08):
    """背景点阵。"""
    step = max(0.2, density * 12)
    x = 0.3
    while x < 13.1:
        y = 0.3
        while y < 7.2:
            dot = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(x), Inches(y), Inches(0.05), Inches(0.05))
            set_shape_fill(dot, color, alpha)
            dot.line.fill.background()
            no_shadow(dot)
            y += step
        x += step


def add_deco_lines(slide, prs, color, alpha=20, count=3):
    """右上角斜线装饰。"""
    import random
    rnd = random.Random(42)
    for i in range(count):
        w = 2.5 + rnd.random() * 2
        ln = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Emu(int(prs.slide_width * (0.78 - i * 0.08))),
            Emu(int(prs.slide_height * (0.10 + i * 0.14))),
            Inches(w), Inches(0.045))
        ln.rotation = 30
        set_shape_fill(ln, color, alpha)
        ln.line.fill.background()
        no_shadow(ln)


def add_deco_blob(slide, prs, color, alpha=14, count=2):
    """大色块圆形装饰。"""
    import random
    rnd = random.Random(7)
    for i in range(count):
        size = 2.2 + rnd.random() * 1.8
        cx = rnd.uniform(-0.8, 11.5)
        cy = rnd.uniform(-1.0, 6.0)
        blob = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(cx), Inches(cy), Inches(size), Inches(size))
        set_shape_fill(blob, color, alpha)
        blob.line.fill.background()
        no_shadow(blob)


def add_deco_corner(slide, prs, color, alpha=100):
    """左上角折角装饰。"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_TRIANGLE, 0, 0, Inches(0.9), Inches(0.9))
    shape.rotation = 180
    set_shape_fill(shape, color, alpha)
    shape.line.fill.background()
    no_shadow(shape)


def add_deco_glow(slide, prs, color, x_frac=0.85, y_frac=0.15, size=3.5, alpha=22):
    """光斑装饰（大圆形柔光）。"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Emu(int(prs.slide_width * x_frac) - Inches(size / 2)),
        Emu(int(prs.slide_height * y_frac) - Inches(size / 2)),
        Inches(size), Inches(size))
    set_shape_fill(shape, color, alpha)
    shape.line.fill.background()
    no_shadow(shape)


def add_deco_grid(slide, prs, color, alpha=14, spacing=0.9):
    """网格背景线。"""
    x = spacing
    while x < prs.slide_width.inches - 0.1:
        vline = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(x), 0, Inches(0.012), prs.slide_height)
        set_shape_fill(vline, color, alpha)
        vline.line.fill.background()
        no_shadow(vline)
        x += spacing
    y = spacing
    while y < prs.slide_height.inches - 0.1:
        hline = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, Inches(y), prs.slide_width, Inches(0.012))
        set_shape_fill(hline, color, alpha)
        hline.line.fill.background()
        no_shadow(hline)
        y += spacing


def add_deco_side_glow(slide, prs, color, side="left", width=4.0, alpha=20):
    """侧边光晕（左侧/右侧大面积柔光带）。"""
    if side == "right":
        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Emu(int(prs.slide_width - Inches(width))),
            Emu(int(prs.slide_height / 2) - Inches(3.2)),
            Inches(width * 2), Inches(6.4))
    else:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Emu(-Inches(width)),
            Emu(int(prs.slide_height / 2) - Inches(3.2)),
            Inches(width * 2), Inches(6.4))
    set_shape_fill(shape, color, alpha)
    shape.line.fill.background()
    no_shadow(shape)


def add_deco_waves(slide, prs, color, alpha=16, count=3):
    """底部波浪线装饰。"""
    import random
    rnd = random.Random(99)
    for i in range(count):
        y_base = prs.slide_height.inches - 0.4 - i * 0.55
        amp = 0.18 + rnd.random() * 0.12
        # 用一系列小圆弧段近似波浪
        seg_w = 0.5
        x = 0.0
        up = True
        while x < prs.slide_width.inches - seg_w:
            arc = slide.shapes.add_shape(
                MSO_SHAPE.CHORD,
                Inches(x), Inches(y_base - amp if up else y_base),
                Inches(seg_w), Inches(amp))
            set_shape_fill(arc, color, alpha)
            arc.line.fill.background()
            no_shadow(arc)
            x += seg_w
            up = not up


def apply_decoration(slide, deco, theme, prs):
    """应用单条装饰。deco: dict。"""
    if deco is None:
        return
    d = dict(deco)
    d.setdefault("color", theme["colors"]["primary"])
    color = color_hex(d.get("color"), theme)
    alpha = int(d.get("alpha", 100))
    dtype = str(d.get("type", "top_bar")).lower()
    if dtype == "top_bar":
        add_deco_bar(slide, prs, color, top=True, height=d.get("height", 0.08), alpha=alpha)
    elif dtype == "bottom_bar":
        add_deco_bar(slide, prs, color, top=False, height=d.get("height", 0.08), alpha=alpha)
    elif dtype == "circle":
        add_deco_circle(slide, prs, color,
                        d.get("x_frac", 0.92), d.get("y_frac", 0.08),
                        d.get("size", 0.16), alpha)
    elif dtype == "dot_grid":
        add_dot_grid(slide, prs, color, alpha, d.get("density", 0.08))
    elif dtype == "lines":
        add_deco_lines(slide, prs, color, alpha, d.get("count", 3))
    elif dtype == "blob":
        add_deco_blob(slide, prs, color, alpha, d.get("count", 2))
    elif dtype == "corner":
        add_deco_corner(slide, prs, color, alpha)
    elif dtype == "glow":
        add_deco_glow(slide, prs, color,
                      d.get("x_frac", 0.85), d.get("y_frac", 0.15),
                      d.get("size", 3.5), alpha)
    elif dtype == "grid":
        add_deco_grid(slide, prs, color, alpha, d.get("spacing", 0.9))
    elif dtype == "side_glow":
        add_deco_side_glow(slide, prs, color, d.get("side", "left"),
                           d.get("width", 4.0), alpha)
    elif dtype == "waves":
        add_deco_waves(slide, prs, color, alpha, d.get("count", 3))


def apply_decorations(slide, decorations, theme, prs):
    """应用多条装饰；空则用主题默认装饰。"""
    if not decorations:
        deco = theme.get("decoration")
        if deco:
            apply_decoration(slide, deco, theme, prs)
        return
    for deco in decorations:
        apply_decoration(slide, deco, theme, prs)


# ------------------------------------------------------------------ 文本估算
def estimate_lines(text, size, width_in, char_w_ratio=0.98):
    """估算文本行数（粗略）。"""
    if not text:
        return 1
    char_w = size * char_w_ratio / 72.0
    per_line = max(1, int(width_in / char_w))
    lines = 0
    for seg in str(text).split("\n"):
        lines += max(1, math.ceil(len(seg) / per_line))
    return lines


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    """创建文本框，返回 text_frame。"""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return tf
