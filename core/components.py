# -*- coding: utf-8 -*-
"""组件库：20+ 可视化组件，全部基于 python-pptx 原生形状渲染。

每个组件函数签名统一：
    render(slide, spec, theme, ctx)
其中 ctx 包含：
    prs            Presentation 对象
    slide_width    幻灯片宽度(EMU)
    slide_height   幻灯片高度(EMU)
    fonts          字体字典 {"heading":.., "body":..}
    style          "dark"|"light"

所有组件均支持通用定位字段 x / y / w / h（英寸，或 "50%"、"fill"）。
"""
import os

from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

from . import utils as U

# 通用字段提取
def _geo(spec, ctx):
    x = U.parse_len(spec.get("x", 0), ctx["slide_width"])
    y = U.parse_len(spec.get("y", 0), ctx["slide_height"])
    w = U.parse_len(spec.get("w", 4), ctx["slide_width"])
    h = U.parse_len(spec.get("h", 1), ctx["slide_height"])
    return x, y, w, h


def _f(spec, key, default, theme, colors_key=None):
    """取组件级颜色覆盖，否则主题色。"""
    if key in spec and spec[key] is not None:
        return U.color_hex(spec[key], theme)
    if colors_key:
        return theme["colors"].get(colors_key, default)
    return U.color_hex(default, theme)


def _font(spec, ctx, key="font"):
    return spec.get(key) or ctx["fonts"].get("body", "微软雅黑")


# ================================================================== 1. title
def render_title(slide, spec, theme, ctx):
    x, y, w, h = _geo(spec, ctx)
    text = spec.get("text", "")
    size = spec.get("size", 32)
    color = _f(spec, "color", "text", theme, "text")
    align = spec.get("align", "left")
    align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
    tf = U.textbox(slide, x, y, w, h)
    p = U.add_para(tf, first=True, align=align_map.get(align, PP_ALIGN.LEFT))
    r = p.add_run()
    U.style_run(r, text=text, size=size, bold=True, color=color, font=_font(spec, ctx, "heading"))
    # 可选的左侧强调条
    if spec.get("accent_bar"):
        bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(0.09), h)
        U.set_shape_fill(bar, _f(spec, "accent", "primary", theme, "primary"))
        U.set_round_rect_radius(bar, 0.5)
        bar.line.fill.background()
        U.no_shadow(bar)


# ================================================================== 2. text
def render_text(slide, spec, theme, ctx):
    x, y, w, h = _geo(spec, ctx)
    tf = U.textbox(slide, x, y, w, h)
    paragraphs = spec.get("paragraphs") or [spec.get("content") or [{"text": spec.get("text", "")}]]
    default_size = spec.get("size", 16)
    default_color = _f(spec, "color", "text", theme, "text")
    align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
    align = align_map.get(spec.get("align", "left"), PP_ALIGN.LEFT)
    line_spacing = spec.get("line_spacing", 1.15)
    for pi, para_spec in enumerate(paragraphs):
        p = U.add_para(tf, first=(pi == 0), align=align, line_spacing=line_spacing,
                       space_after=spec.get("space_after"))
        if isinstance(para_spec, str):
            para_spec = [{"text": para_spec}]
        for run_spec in para_spec:
            r = p.add_run()
            rs = dict(run_spec) if isinstance(run_spec, dict) else {"text": str(run_spec)}
            U.style_run(
                r,
                text=rs.get("text", ""),
                size=rs.get("size", default_size),
                bold=rs.get("bold", False),
                italic=rs.get("italic", False),
                color=U.color_hex(rs.get("color", default_color), theme),
                font=rs.get("font") or _font(spec, ctx),
            )


# ================================================================== 3. list
def render_list(slide, spec, theme, ctx):
    x, y, w, h = _geo(spec, ctx)
    tf = U.textbox(slide, x, y, w, h)
    items = spec.get("items", [])
    marker = spec.get("marker", "bullet")  # bullet | number | icon
    size = spec.get("size", 17)
    color = _f(spec, "color", "text", theme, "text")
    marker_color = _f(spec, "marker_color", "primary", theme, "primary")
    line_spacing = spec.get("line_spacing", 1.3)
    space_after = spec.get("space_after", 8)
    for i, item in enumerate(items):
        if isinstance(item, str):
            item = {"text": item}
        level = int(item.get("level", 0))
        p = U.add_para(tf, first=(i == 0), line_spacing=line_spacing,
                       space_after=space_after, indent_level=0)
        indent = "　" * level
        # 标记
        if marker == "number":
            mark = f"{i + 1}. "
            p.alignment = PP_ALIGN.LEFT
            r = p.add_run()
            U.style_run(r, text=indent + mark, size=size, bold=True, color=marker_color,
                        font=_font(spec, ctx))
        elif marker == "icon":
            mark = item.get("icon", "▸") + " "
            r = p.add_run()
            U.style_run(r, text=indent + mark, size=size, bold=True, color=marker_color,
                        font=_font(spec, ctx))
        else:
            mark = "• "
            r = p.add_run()
            U.style_run(r, text=indent + mark, size=size, bold=True, color=marker_color,
                        font=_font(spec, ctx))
        # 内容
        r2 = p.add_run()
        U.style_run(r2, text=item.get("text", ""), size=size,
                    bold=item.get("bold", False), color=color, font=_font(spec, ctx))


# ================================================================== 4. card
def render_card(slide, spec, theme, ctx):
    x, y, w, h = _geo(spec, ctx)
    accent = _f(spec, "accent", "primary", theme, "primary")
    bg = _f(spec, "bg", "card_bg", theme, "card_bg")
    border = _f(spec, "border", "card_border", theme, "card_border")
    radius = spec.get("radius", 0.08)
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    U.set_shape_fill(card, bg)
    U.set_round_rect_radius(card, radius)
    U.set_shape_line(card, border, spec.get("border_width", 1))
    if spec.get("shadow", True):
        U.add_shadow(card, blur=0.12, dist=0.03, alpha=30)
    else:
        U.no_shadow(card)
    # 顶部强调条
    if spec.get("accent_bar", True):
        bar_w = Inches(0.5)
        bar = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.3), y + Inches(0.3),
            bar_w, Inches(0.06))
        U.set_shape_fill(bar, accent)
        U.set_round_rect_radius(bar, 0.5)
        bar.line.fill.background()
        U.no_shadow(bar)
    pad = Inches(0.3)
    inner_x = x + pad
    inner_y = y + pad
    inner_w = w - 2 * pad
    # 图标
    icon = spec.get("icon")
    icon_top = y + Inches(0.28)
    if icon:
        icon_size = spec.get("icon_size", 0.6)
        ic = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, inner_x, icon_top, Inches(icon_size), Inches(icon_size))
        U.set_shape_fill(ic, accent, 18)
        ic.line.fill.background()
        U.no_shadow(ic)
        itf = ic.text_frame
        itf.word_wrap = False
        itf.vertical_anchor = MSO_ANCHOR.MIDDLE
        itf.margin_left = itf.margin_right = itf.margin_top = itf.margin_bottom = 0
        ip = itf.paragraphs[0]
        ip.alignment = PP_ALIGN.CENTER
        ir = ip.add_run()
        U.style_run(ir, text=icon, size=int(icon_size * 30), font=_font(spec, ctx))
        text_x = inner_x + Inches(icon_size) + Inches(0.15)
        text_w = inner_w - Inches(icon_size) - Inches(0.15)
    else:
        text_x = inner_x
        text_w = inner_w
    # 标题 + 文本
    tf = U.textbox(slide, text_x, icon_top if icon else inner_y, text_w, h - 2 * pad)
    title = spec.get("title")
    if title:
        p = U.add_para(tf, first=True)
        r = p.add_run()
        U.style_run(r, text=title, size=spec.get("title_size", 18), bold=True,
                    color=_f(spec, "title_color", "text", theme, "text"),
                    font=_font(spec, ctx, "heading"))
    text = spec.get("text")
    if text:
        p2 = U.add_para(tf, first=(not title), space_before=6, line_spacing=1.2)
        r2 = p2.add_run()
        U.style_run(r2, text=text, size=spec.get("text_size", 13),
                    color=_f(spec, "text_color", "text_muted", theme, "text_muted"),
                    font=_font(spec, ctx))


# ================================================================== 5. kpi
def render_kpi(slide, spec, theme, ctx):
    x, y, w, h = _geo(spec, ctx)
    accent = _f(spec, "accent", "primary", theme, "primary")
    bg = _f(spec, "bg", "card_bg", theme, "card_bg")
    border = _f(spec, "border", "card_border", theme, "card_border")
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    U.set_shape_fill(card, bg)
    U.set_round_rect_radius(card, spec.get("radius", 0.1))
    U.set_shape_line(card, border, 1)
    U.add_shadow(card, blur=0.14, dist=0.035, alpha=32)
    pad = Inches(0.3)
    tf = U.textbox(slide, x + pad, y + pad, w - 2 * pad, h - 2 * pad,
                   anchor=MSO_ANCHOR.MIDDLE)
    # 图标行
    if spec.get("icon"):
        ip = U.add_para(tf, first=True)
        ir = ip.add_run()
        U.style_run(ir, text=spec["icon"], size=24, font=_font(spec, ctx))
    # 数值
    vp = U.add_para(tf, first=(not spec.get("icon")), space_before=4)
    vr = vp.add_run()
    U.style_run(vr, text=spec.get("value", ""), size=spec.get("value_size", 34),
                bold=True, color=accent, font=_font(spec, ctx, "heading"))
    # 单位
    if spec.get("unit"):
        ur = vp.add_run()
        U.style_run(ur, text=" " + spec["unit"], size=spec.get("unit_size", 15),
                    color=_f(spec, "unit_color", "text_muted", theme, "text_muted"),
                    font=_font(spec, ctx))
    # 标签
    lp = U.add_para(tf, space_before=4)
    lr = lp.add_run()
    U.style_run(lr, text=spec.get("label", ""), size=13,
                color=_f(spec, "label_color", "text_muted", theme, "text_muted"),
                font=_font(spec, ctx))
    # 增量徽标
    if spec.get("delta"):
        dp = U.add_para(tf, space_before=6)
        dr = dp.add_run()
        delta = str(spec["delta"])
        up = not delta.startswith("-")
        U.style_run(dr, text=("▲ " if up else "▼ ") + delta, size=12, bold=True,
                    color=theme["colors"]["success" if up else "danger"],
                    font=_font(spec, ctx))


# ================================================================== 6. badge
def render_badge(slide, spec, theme, ctx):
    x, y, w, h = _geo(spec, ctx)
    text = spec.get("text", "")
    variant = spec.get("variant", "fill")  # fill | outline | soft
    color = _f(spec, "color", "primary", theme, "primary")
    size = spec.get("size", 14)
    # 估算宽度
    est_w = max(0.7, len(text) * size * 0.16 + 0.5)
    if w == Inches(4):  # 默认宽度时用估算
        w = Inches(est_w)
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    U.set_round_rect_radius(badge, 0.5)
    if variant == "outline":
        U.set_shape_fill(badge, None)
        U.set_shape_line(badge, color, 1.5)
        txt_color = color
    elif variant == "soft":
        U.set_shape_fill(badge, color, 18)
        badge.line.fill.background()
        txt_color = color
    else:
        U.set_shape_fill(badge, color)
        badge.line.fill.background()
        txt_color = theme["colors"].get("text_invert", "FFFFFF")
    U.no_shadow(badge)
    tf = badge.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.1)
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    U.style_run(r, text=text, size=size, bold=True, color=txt_color, font=_font(spec, ctx))


# ================================================================== 7. divider
def render_divider(slide, spec, theme, ctx):
    x, y, w, h = _geo(spec, ctx)
    color = _f(spec, "color", "primary", theme, "primary")
    style = spec.get("style", "line")  # line | gradient | dot
    thickness = spec.get("thickness", 0.025)
    if style == "dot":
        count = spec.get("count", 12)
        dot_w = spec.get("dot_size", 0.07)
        gap = (w - count * Inches(dot_w)) / (count - 1) if count > 1 else 0
        for i in range(count):
            d = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                x + Emu(int(i * (Inches(dot_w) + gap))), y,
                Inches(dot_w), Inches(dot_w))
            U.set_shape_fill(d, color, 100 if i < spec.get("filled", count) else 30)
            d.line.fill.background()
            U.no_shadow(d)
    elif style == "gradient":
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Inches(thickness))
        U.make_gradient_fill(bar, [color, _f(spec, "color2", "secondary", theme, "secondary")], angle=0)
        bar.line.fill.background()
        U.no_shadow(bar)
    else:
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Inches(thickness))
        U.set_shape_fill(bar, color, spec.get("alpha", 60))
        bar.line.fill.background()
        U.no_shadow(bar)


# ================================================================== 8. shape
def render_shape(slide, spec, theme, ctx):
    x, y, w, h = _geo(spec, ctx)
    shape_map = {
        "rect": MSO_SHAPE.RECTANGLE,
        "round_rect": MSO_SHAPE.ROUNDED_RECTANGLE,
        "circle": MSO_SHAPE.OVAL,
        "oval": MSO_SHAPE.OVAL,
        "line": MSO_SHAPE.RECTANGLE,
        "right_arrow": MSO_SHAPE.RIGHT_ARROW,
        "chevron": MSO_SHAPE.CHEVRON,
        "pentagon": MSO_SHAPE.PENTAGON,
        "hexagon": MSO_SHAPE.HEXAGON,
        "triangle": MSO_SHAPE.ISOCELES_TRIANGLE,
        "diamond": MSO_SHAPE.DIAMOND,
        "star": MSO_SHAPE.STAR_5_POINT,
    }
    stype = str(spec.get("shape", "rect")).lower()
    mso = shape_map.get(stype, MSO_SHAPE.RECTANGLE)
    shape = slide.shapes.add_shape(mso, x, y, w, h)
    fill = spec.get("fill")
    if fill == "none":
        shape.fill.background()
    elif fill == "gradient":
        colors = [U.color_hex(c, theme) for c in spec.get("colors", ["primary", "secondary"])]
        U.make_gradient_fill(shape, colors, spec.get("angle", 90))
    else:
        U.set_shape_fill(shape, _f(spec, "fill_color", "primary", theme, "primary"),
                         spec.get("alpha"))
    line = spec.get("line")
    if line == "none":
        shape.line.fill.background()
    else:
        U.set_shape_line(shape, _f(spec, "line_color", "card_border", theme, "card_border"),
                         spec.get("line_width", 1))
    if stype == "round_rect":
        U.set_round_rect_radius(shape, spec.get("radius", 0.15))
    if spec.get("rotation"):
        shape.rotation = float(spec["rotation"])
    if spec.get("shadow", False):
        U.add_shadow(shape, alpha=30)
    else:
        U.no_shadow(shape)
    # 文本
    text = spec.get("text")
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = tf.margin_right = Inches(0.05)
        tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        U.style_run(r, text=text, size=spec.get("text_size", 16), bold=True,
                    color=_f(spec, "text_color", "text_invert", theme, "text_invert"),
                    font=_font(spec, ctx))


# ================================================================== 9. icon
def render_icon(slide, spec, theme, ctx):
    x, y = U.parse_len(spec.get("x", 0), ctx["slide_width"]), U.parse_len(spec.get("y", 0), ctx["slide_height"])
    size = Inches(spec.get("size", 0.7))
    emoji = spec.get("emoji", spec.get("text", "✦"))
    color = _f(spec, "color", "primary", theme, "primary")
    bg_alpha = spec.get("bg_alpha", 18)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, size, size)
    if spec.get("bg", True):
        U.set_shape_fill(circle, color, bg_alpha)
    else:
        circle.fill.background()
    circle.line.fill.background()
    U.no_shadow(circle)
    tf = circle.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    U.style_run(r, text=emoji, size=int(spec.get("emoji_size", size.inches * 28)),
                font=_font(spec, ctx))


# ================================================================== 10. progress
def render_progress(slide, spec, theme, ctx):
    x, y, w, h = _geo(spec, ctx)
    value = max(0.0, min(1.0, float(spec.get("value", 0.5))))
    color = _f(spec, "color", "primary", theme, "primary")
    track = _f(spec, "track", "card_border", theme, "card_border")
    bar_h = Inches(spec.get("bar_height", 0.22))
    # 轨道
    track_sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, bar_h)
    U.set_shape_fill(track_sh, track, 40)
    U.set_round_rect_radius(track_sh, 0.5)
    track_sh.line.fill.background()
    U.no_shadow(track_sh)
    # 前景
    fw = Emu(int(w * value))
    if fw > 0:
        fg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, fw, bar_h)
        if spec.get("gradient"):
            U.make_gradient_fill(fg, [color, _f(spec, "color2", "secondary", theme, "secondary")], angle=0)
        else:
            U.set_shape_fill(fg, color)
        U.set_round_rect_radius(fg, 0.5)
        fg.line.fill.background()
        U.no_shadow(fg)
    # 标签（右侧百分比 或 左侧 label）
    label = spec.get("label")
    tf = U.textbox(slide, x + w + Inches(0.15), y - Inches(0.05),
                   Inches(1.5), Inches(0.3))
    p = U.add_para(tf, first=True)
    r = p.add_run()
    show = label if label else f"{int(value * 100)}%"
    U.style_run(r, text=show, size=12, bold=True, color=color, font=_font(spec, ctx))


# ================================================================== 11. bar_chart
def render_bar_chart(slide, spec, theme, ctx):
    x, y, w, h = _geo(spec, ctx)
    data = spec.get("data", [])
    if not data:
        return
    color = _f(spec, "color", "primary", theme, "primary")
    color2 = _f(spec, "color2", "secondary", theme, "secondary")
    bar_h = Inches(spec.get("bar_height", 0.4))
    gap = Inches(spec.get("gap", 0.35))
    label_w = Inches(spec.get("label_width", 1.4))
    value_w = Inches(spec.get("value_width", 1.0))
    chart_w = max(w - label_w - value_w, Inches(0.2))
    max_val = max([float(d.get("value", 0)) for d in data] + [1.0])
    bar_x = x + label_w
    label_tf = U.textbox(slide, x, y, label_w, h)
    value_tf = U.textbox(slide, x + w - value_w, y, value_w, h)
    for i, d in enumerate(data):
        by = y + Emu(int(i * (bar_h + gap)))
        # 标签
        lp = U.add_para(label_tf, first=(i == 0))
        lr = lp.add_run()
        U.style_run(lr, text=str(d.get("label", "")), size=13, bold=True,
                    color=_f(spec, "label_color", "text", theme, "text"),
                    font=_font(spec, ctx))
        # 值
        vp = U.add_para(value_tf, first=(i == 0), align=PP_ALIGN.RIGHT)
        vr = vp.add_run()
        U.style_run(vr, text=str(d.get("value", "")), size=13, bold=True,
                    color=color, font=_font(spec, ctx))
        # 条形
        frac = float(d.get("value", 0)) / max_val
        bw = Emu(max(int(chart_w * frac), Inches(0.05)))
        bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bar_x, by, bw, bar_h)
        if spec.get("gradient") and i % 2 == 0:
            U.make_gradient_fill(bar, [color, color2], angle=0)
        else:
            U.set_shape_fill(bar, color if i % 2 == 0 else color2,
                             spec.get("alpha"))
        U.set_round_rect_radius(bar, 0.5)
        bar.line.fill.background()
        U.no_shadow(bar)


# ================================================================== 12. pie_chart
def render_pie_chart(slide, spec, theme, ctx):
    x, y = U.parse_len(spec.get("x", 0), ctx["slide_width"]), U.parse_len(spec.get("y", 0), ctx["slide_height"])
    size = Inches(spec.get("size", 3.2))
    data = spec.get("data", [])
    if not data:
        return
    total = sum(float(d.get("value", 0)) for d in data) or 1.0
    donut = spec.get("donut", False)
    # 预定义色板（按主题色循环）
    palette = [theme["colors"].get(k) for k in
               ("primary", "secondary", "accent", "success", "warning", "danger")]
    start = 0.0
    for i, d in enumerate(data):
        frac = float(d.get("value", 0)) / total
        if frac <= 0:
            continue
        color = U.color_hex(d.get("color", palette[i % len(palette)]), theme)
        pie = slide.shapes.add_shape(MSO_SHAPE.PIE, x, y, size, size)
        pie.adjustments[0] = start
        pie.adjustments[1] = start + frac
        U.set_shape_fill(pie, color)
        pie.line.color.rgb = U.hex_to_rgb(theme["colors"]["card_bg"])
        pie.line.width = Pt(1.5)
        U.no_shadow(pie)
        start += frac
    if donut:
        # 中心孔
        hole_size = size * 0.55
        hole = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            x + Emu(int((size - hole_size) / 2)),
            y + Emu(int((size - hole_size) / 2)),
            hole_size, hole_size)
        hole.fill.solid()
        hole.fill.fore_color.rgb = U.hex_to_rgb(
            _f(spec, "hole_color", "card_bg", theme, "card_bg"))
        hole.line.fill.background()
        U.no_shadow(hole)
        if spec.get("center_text"):
            htf = hole.text_frame
            htf.word_wrap = True
            htf.vertical_anchor = MSO_ANCHOR.MIDDLE
            htf.margin_left = htf.margin_right = Inches(0.05)
            htf.margin_top = htf.margin_bottom = 0
            hp = htf.paragraphs[0]
            hp.alignment = PP_ALIGN.CENTER
            hr = hp.add_run()
            U.style_run(hr, text=spec["center_text"], size=16, bold=True,
                        color=_f(spec, "center_color", "primary", theme, "primary"),
                        font=_font(spec, ctx, "heading"))
    # 图例
    if spec.get("legend", True):
        lx = x + size + Inches(0.5)
        ly = y
        legend_w = Inches(3.2)
        ltf = U.textbox(slide, lx, ly, legend_w, size)
        for i, d in enumerate(data):
            color = U.color_hex(d.get("color", palette[i % len(palette)]), theme)
            p = U.add_para(ltf, first=(i == 0), space_after=6)
            dot = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, lx, ly + Emu(int(i * (0.26 * 914400))), Inches(0.14), Inches(0.14))
            U.set_shape_fill(dot, color)
            dot.line.fill.background()
            U.no_shadow(dot)
            r = p.add_run()
            U.style_run(r, text=f"{d.get('label','')}  {int(float(d.get('value',0)))}",
                        size=13, color=_f(spec, "legend_color", "text", theme, "text"),
                        font=_font(spec, ctx))


# ================================================================== 13. timeline
def render_timeline(slide, spec, theme, ctx):
    x, y, w, h = _geo(spec, ctx)
    items = spec.get("items", [])
    if not items:
        return
    color = _f(spec, "color", "primary", theme, "primary")
    vertical = spec.get("vertical", False)
    if vertical:
        # 纵向：左侧竖线 + 节点 + 右侧卡片
        line_x = x + Inches(0.25)
        line_top = y
        line_bottom = y + h
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      line_x, line_top, Inches(0.04), h)
        U.set_shape_fill(line, color, 40)
        line.line.fill.background()
        U.no_shadow(line)
        card_w = w - Inches(0.9)
        node_size = Inches(0.24)
        step = h / max(1, len(items))
        for i, item in enumerate(items):
            ny = y + Emu(int((i + 0.5) * step))
            node = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                line_x - Inches(0.1), ny - Inches(0.12), node_size, node_size)
            U.set_shape_fill(node, color)
            node.line.color.rgb = U.hex_to_rgb(_f(spec, "bg", "card_bg", theme, "card_bg"))
            node.line.width = Pt(2)
            U.no_shadow(node)
            # 卡片
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                line_x + Inches(0.5), ny - Inches(0.35), card_w, Inches(0.7))
            U.set_shape_fill(card, _f(spec, "bg", "card_bg", theme, "card_bg"))
            U.set_round_rect_radius(card, 0.12)
            U.set_shape_line(card, _f(spec, "border", "card_border", theme, "card_border"), 1)
            U.no_shadow(card)
            ctf = card.text_frame
            ctf.word_wrap = True
            ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
            ctf.margin_left = Inches(0.2)
            ctf.margin_right = Inches(0.15)
            ctf.margin_top = ctf.margin_bottom = Inches(0.05)
            cp = ctf.paragraphs[0]
            if item.get("date"):
                cr = cp.add_run()
                U.style_run(cr, text=item["date"] + "  ", size=12, bold=True,
                            color=color, font=_font(spec, ctx))
            cr2 = cp.add_run()
            U.style_run(cr2, text=item.get("title", ""), size=14, bold=True,
                        color=_f(spec, "title_color", "text", theme, "text"),
                        font=_font(spec, ctx, "heading"))
            if item.get("desc"):
                dp2 = ctf.add_paragraph()
                dr2 = dp2.add_run()
                U.style_run(dr2, text=item["desc"], size=11,
                            color=_f(spec, "desc_color", "text_muted", theme, "text_muted"),
                            font=_font(spec, ctx))
    else:
        # 横向：一条主线 + 各节点上下错落卡片
        line_y = y + Inches(1.1)
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      x, line_y, w, Inches(0.04))
        U.set_shape_fill(line, color, 40)
        line.line.fill.background()
        U.no_shadow(line)
        n = len(items)
        slot_w = w / n
        for i, item in enumerate(items):
            cx = x + Emu(int((i + 0.5) * slot_w))
            node = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx - Inches(0.12), line_y - Inches(0.1),
                                          Inches(0.24), Inches(0.24))
            U.set_shape_fill(node, color)
            node.line.fill.background()
            U.no_shadow(node)
            up = i % 2 == 0
            card_h = Inches(0.85)
            cy = line_y - Inches(0.5) - card_h if up else line_y + Inches(0.5)
            card_w = min(slot_w, Inches(2.4))
            card_x = cx - Emu(int(card_w / 2))
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, card_x, cy, card_w, card_h)
            U.set_shape_fill(card, _f(spec, "bg", "card_bg", theme, "card_bg"))
            U.set_round_rect_radius(card, 0.1)
            U.set_shape_line(card, _f(spec, "border", "card_border", theme, "card_border"), 1)
            U.no_shadow(card)
            ctf = card.text_frame
            ctf.word_wrap = True
            ctf.margin_left = Inches(0.15)
            ctf.margin_right = Inches(0.1)
            ctf.margin_top = ctf.margin_bottom = Inches(0.08)
            cp = ctf.paragraphs[0]
            cp.alignment = PP_ALIGN.CENTER
            if item.get("date"):
                cr = cp.add_run()
                U.style_run(cr, text=item["date"] + "\n", size=11, bold=True,
                            color=color, font=_font(spec, ctx))
            cr2 = cp.add_run()
            U.style_run(cr2, text=item.get("title", ""), size=13, bold=True,
                        color=_f(spec, "title_color", "text", theme, "text"),
                        font=_font(spec, ctx, "heading"))
            if item.get("desc"):
                dp2 = ctf.add_paragraph()
                dp2.alignment = PP_ALIGN.CENTER
                dr2 = dp2.add_run()
                U.style_run(dr2, text=item["desc"], size=10,
                            color=_f(spec, "desc_color", "text_muted", theme, "text_muted"),
                            font=_font(spec, ctx))


# ================================================================== 14. steps
def render_steps(slide, spec, theme, ctx):
    x, y, w, h = _geo(spec, ctx)
    items = spec.get("items", [])
    if not items:
        return
    color = _f(spec, "color", "primary", theme, "primary")
    n = len(items)
    step_w = w / n
    circle_d = Inches(0.7)
    current = int(spec.get("current", -1))
    for i, item in enumerate(items):
        cx = x + Emu(int((i + 0.5) * step_w))
        # 连接线
        if i < n - 1:
            lx = cx + Emu(int(circle_d / 2)) + Inches(0.05)
            lw = step_w - circle_d - Inches(0.1)
            line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, lx, y + Inches(0.32), lw, Inches(0.04))
            done = current < 0 or i < current
            U.set_shape_fill(line, color if done else _f(spec, "muted_color", "card_border", theme, "card_border"),
                             spec.get("line_alpha", 50))
            line.line.fill.background()
            U.no_shadow(line)
        # 编号圆
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, cx - Inches(0.35), y, circle_d, circle_d)
        if current >= 0 and i == current:
            U.make_gradient_fill(circle, [color, _f(spec, "color2", "secondary", theme, "secondary")], 90)
            circle.line.fill.background()
            U.add_shadow(circle, alpha=35)
        elif current >= 0 and i < current:
            U.set_shape_fill(circle, color, 60)
            circle.line.fill.background()
            U.no_shadow(circle)
        else:
            U.set_shape_fill(circle, _f(spec, "bg", "card_bg", theme, "card_bg"))
            U.set_shape_line(circle, _f(spec, "border", "card_border", theme, "card_border"), 1.5)
            U.no_shadow(circle)
        ctf = circle.text_frame
        ctf.word_wrap = False
        ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
        ctf.margin_left = ctf.margin_right = ctf.margin_top = ctf.margin_bottom = 0
        cp = ctf.paragraphs[0]
        cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run()
        num_color = theme["colors"]["text_invert"] if (current >= 0 and i <= current) else _f(spec, "num_color", "text_muted", theme, "text_muted")
        U.style_run(cr, text=str(i + 1), size=20, bold=True, color=num_color, font=_font(spec, ctx, "heading"))
        # 标题
        label = item if isinstance(item, str) else item.get("label", "")
        lt = U.textbox(slide, cx - Inches(1.6), y + circle_d + Inches(0.15), Inches(3.2), Inches(0.6))
        lp = U.add_para(lt, first=True, align=PP_ALIGN.CENTER)
        lr = lp.add_run()
        U.style_run(lr, text=label, size=14, bold=True,
                    color=_f(spec, "label_color", "text", theme, "text"),
                    font=_font(spec, ctx, "heading"))
        if isinstance(item, dict) and item.get("desc"):
            dt = U.textbox(slide, cx - Inches(1.7), y + circle_d + Inches(0.65), Inches(3.4), Inches(0.5))
            dp = U.add_para(dt, first=True, align=PP_ALIGN.CENTER)
            dr = dp.add_run()
            U.style_run(dr, text=item["desc"], size=11,
                        color=_f(spec, "desc_color", "text_muted", theme, "text_muted"),
                        font=_font(spec, ctx))


# ================================================================== 15. grid_cards
def render_grid_cards(slide, spec, theme, ctx):
    x, y, w, h = _geo(spec, ctx)
    cards = spec.get("cards", [])
    if not cards:
        return
    cols = int(spec.get("cols", 3))
    gap = Inches(spec.get("gap", 0.35))
    card_w = Emu(int((w - gap * (cols - 1)) / cols))
    rows = (len(cards) + cols - 1) // cols
    card_h = spec.get("card_height", 2.0)
    if isinstance(card_h, (int, float)) and card_h <= 10:
        card_h = Inches(card_h)
    for i, c in enumerate(cards):
        r_i, c_i = divmod(i, cols)
        cx = x + Emu(int(c_i * (card_w + gap)))
        cy = y + Emu(int(r_i * (card_h + gap)))
        sub = dict(c)
        sub.setdefault("x", cx)
        sub.setdefault("y", cy)
        sub.setdefault("w", card_w)
        sub.setdefault("h", card_h)
        sub["shadow"] = spec.get("shadow", True)
        render_card(slide, sub, theme, ctx)


# ================================================================== 16. table
def render_table(slide, spec, theme, ctx):
    x, y, w, h = _geo(spec, ctx)
    headers = spec.get("headers", [])
    rows = spec.get("rows", [])
    if not headers:
        return
    n_rows = len(rows) + 1
    n_cols = len(headers)
    header_bg = _f(spec, "header_bg", "primary", theme, "primary")
    header_fg = _f(spec, "header_fg", "text_invert", theme, "text_invert")
    row_h = spec.get("row_height", 0.5)
    if isinstance(row_h, (int, float)) and row_h <= 10:
        row_h = Inches(row_h)
    table_shape = slide.shapes.add_table(
        n_rows, n_cols, x, y, w, Emu(int(row_h * n_rows)))
    table = table_shape.table
    # 列宽
    col_w = Emu(int(w / n_cols))
    for c in range(n_cols):
        table.columns[c].width = col_w
    # 表头
    for c, ht in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = str(ht)
        cell.fill.solid()
        cell.fill.fore_color.rgb = U.hex_to_rgb(header_bg)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                U.style_run(r, size=spec.get("header_size", 15), bold=True,
                            color=header_fg, font=_font(spec, ctx, "heading"))
    # 数据行
    row_bg = _f(spec, "row_bg", "card_bg", theme, "card_bg")
    row_bg2 = _f(spec, "row_bg2", "card_bg2", theme, "card_bg2")
    row_fg = _f(spec, "row_fg", "text", theme, "text")
    first_col_emphasis = spec.get("first_col_emphasis", False)
    for ri, row in enumerate(rows, start=1):
        bg = row_bg if ri % 2 == 1 else row_bg2
        for ci, val in enumerate(row):
            if ci >= n_cols:
                break
            cell = table.cell(ri, ci)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = U.hex_to_rgb(bg)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            align = PP_ALIGN.CENTER if (ci == 0 or spec.get("center", False)) else PP_ALIGN.LEFT
            for p in cell.text_frame.paragraphs:
                p.alignment = align
                for r in p.runs:
                    bold = first_col_emphasis and ci == 0
                    U.style_run(r, size=spec.get("body_size", 13), bold=bold,
                                color=(_f(spec, "first_col_fg", "primary", theme, "primary")
                                       if (bold and spec.get("first_col_color", True)) else row_fg),
                                font=_font(spec, ctx))
    # 表头下方强调线
    if spec.get("header_line", True):
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x, y + Emu(int(row_h)), w, Inches(0.03))
        U.set_shape_fill(line, header_bg)
        line.line.fill.background()
        U.no_shadow(line)


# ================================================================== 17. quote
def render_quote(slide, spec, theme, ctx):
    x, y, w, h = _geo(spec, ctx)
    text = spec.get("text", "")
    author = spec.get("author")
    accent = _f(spec, "accent", "primary", theme, "primary")
    bg = _f(spec, "bg", "card_bg", theme, "card_bg")
    # 背景卡片
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    U.set_shape_fill(card, bg)
    U.set_round_rect_radius(card, spec.get("radius", 0.08))
    U.set_shape_line(card, _f(spec, "border", "card_border", theme, "card_border"), 1)
    U.no_shadow(card)
    # 左侧竖条
    bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y + Inches(0.2),
                                 Inches(0.08), h - Inches(0.4))
    U.set_shape_fill(bar, accent)
    U.set_round_rect_radius(bar, 0.5)
    bar.line.fill.background()
    U.no_shadow(bar)
    # 引号
    qx = x + Inches(0.4)
    qw = w - Inches(0.8)
    tf = U.textbox(slide, qx, y + Inches(0.25), qw, h - Inches(0.5),
                   anchor=MSO_ANCHOR.MIDDLE)
    p = U.add_para(tf, first=True, line_spacing=1.3)
    r = p.add_run()
    U.style_run(r, text="“", size=spec.get("quote_size", 30), bold=True, color=accent,
                font=_font(spec, ctx, "heading"))
    r2 = p.add_run()
    U.style_run(r2, text=text, size=spec.get("size", 17),
                color=_f(spec, "color", "text", theme, "text"), font=_font(spec, ctx))
    r3 = p.add_run()
    U.style_run(r3, text="”", size=spec.get("quote_size", 30), bold=True, color=accent,
                font=_font(spec, ctx, "heading"))
    if author:
        ap = U.add_para(tf, space_before=10, align=PP_ALIGN.RIGHT)
        ar = ap.add_run()
        U.style_run(ar, text="—— " + author, size=12,
                    color=_f(spec, "author_color", "text_muted", theme, "text_muted"),
                    font=_font(spec, ctx))


# ================================================================== 18. image
def render_image(slide, spec, theme, ctx):
    x, y, w, h = _geo(spec, ctx)
    path = spec.get("path") or spec.get("image_path")
    if not path:
        return
    try:
        from pptx.util import Inches as _In
        pic = slide.shapes.add_picture(path, x, y, width=w)
        if spec.get("h"):
            # 若指定高度则按高度约束
            pic.height = h
        if spec.get("radius"):
            U.set_rounded_picture(pic, spec["radius"])
        if spec.get("border"):
            pic.line.color.rgb = U.hex_to_rgb(_f(spec, "border_color", "card_border", theme, "card_border"))
            pic.line.width = Pt(spec.get("border_width", 2))
        if spec.get("shadow", False):
            U.add_shadow(pic, alpha=30)
    except Exception:
        pass


# ================================================================== 19. team
def render_team(slide, spec, theme, ctx):
    x, y, w, h = _geo(spec, ctx)
    members = spec.get("members", [])
    if not members:
        return
    n = len(members)
    slot_w = w / n
    avatar = Inches(0.9)
    for i, m in enumerate(members):
        cx = x + Emu(int((i + 0.5) * slot_w))
        color = U.color_hex(m.get("color", theme["colors"]["primary"]), theme)
        # 头像圆
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx - avatar / 2, y, avatar, avatar)
        U.set_shape_fill(circle, color, m.get("bg_alpha", 20))
        circle.line.color.rgb = U.hex_to_rgb(color)
        circle.line.width = Pt(2)
        U.no_shadow(circle)
        ctf = circle.text_frame
        ctf.word_wrap = False
        ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
        ctf.margin_left = ctf.margin_right = ctf.margin_top = ctf.margin_bottom = 0
        cp = ctf.paragraphs[0]
        cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run()
        U.style_run(cr, text=m.get("avatar", "🙂"), size=int(avatar.inches * 26),
                    font=_font(spec, ctx))
        # 名字
        nt = U.textbox(slide, cx - Inches(1.4), y + avatar + Inches(0.15), Inches(2.8), Inches(0.4))
        np_ = U.add_para(nt, first=True, align=PP_ALIGN.CENTER)
        nr = np_.add_run()
        U.style_run(nr, text=m.get("name", ""), size=15, bold=True,
                    color=_f(spec, "name_color", "text", theme, "text"),
                    font=_font(spec, ctx, "heading"))
        # 职位
        rt = U.textbox(slide, cx - Inches(1.4), y + avatar + Inches(0.5), Inches(2.8), Inches(0.35))
        rp = U.add_para(rt, first=True, align=PP_ALIGN.CENTER)
        rr = rp.add_run()
        U.style_run(rr, text=m.get("role", ""), size=11,
                    color=_f(spec, "role_color", "text_muted", theme, "text_muted"),
                    font=_font(spec, ctx))


# ================================================================== 20. stat_compare
def render_stat_compare(slide, spec, theme, ctx):
    x, y, w, h = _geo(spec, ctx)
    left = spec.get("left", {})
    right = spec.get("right", {})
    middle = spec.get("middle", {})
    col_w = Emu(int(w / 2))
    bg = _f(spec, "bg", "card_bg", theme, "card_bg")
    for i, side in enumerate([left, right]):
        sx = x + Emu(int(i * col_w))
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, sx, y, col_w, h)
        U.set_shape_fill(card, bg)
        U.set_round_rect_radius(card, 0.08)
        U.set_shape_line(card, _f(spec, "border", "card_border", theme, "card_border"), 1)
        U.no_shadow(card)
        color = U.color_hex(side.get("color", theme["colors"]["primary"] if i == 0 else theme["colors"]["secondary"]), theme)
        tf = U.textbox(slide, sx + Inches(0.35), y + Inches(0.25), col_w - Inches(0.7), h - Inches(0.5),
                       anchor=MSO_ANCHOR.MIDDLE)
        p0 = U.add_para(tf, first=True)
        r0 = p0.add_run()
        U.style_run(r0, text=side.get("label", ""), size=13,
                    color=_f(spec, "label_color", "text_muted", theme, "text_muted"),
                    font=_font(spec, ctx))
        p1 = U.add_para(tf, space_before=4)
        r1 = p1.add_run()
        U.style_run(r1, text=side.get("value", ""), size=spec.get("value_size", 34),
                    bold=True, color=color, font=_font(spec, ctx, "heading"))
        if side.get("desc"):
            p2 = U.add_para(tf, space_before=4)
            r2 = p2.add_run()
            U.style_run(r2, text=side["desc"], size=12,
                        color=_f(spec, "desc_color", "text_muted", theme, "text_muted"),
                        font=_font(spec, ctx))
    # 中间 VS 徽章
    vs = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        x + Emu(int(col_w / 2)) - Inches(0.45), y + Inches(0.2), Inches(0.9), Inches(0.9))
    U.set_shape_fill(vs, _f(spec, "vs_bg", "primary", theme, "primary"))
    vs.line.fill.background()
    U.add_shadow(vs, alpha=35)
    vtf = vs.text_frame
    vtf.word_wrap = False
    vtf.vertical_anchor = MSO_ANCHOR.MIDDLE
    vtf.margin_left = vtf.margin_right = vtf.margin_top = vtf.margin_bottom = 0
    vp = vtf.paragraphs[0]
    vp.alignment = PP_ALIGN.CENTER
    vr = vp.add_run()
    U.style_run(vr, text=middle.get("label", "VS"), size=14, bold=True,
                color=theme["colors"]["text_invert"], font=_font(spec, ctx, "heading"))
    # 中间结果徽标
    if middle.get("value"):
        mt = U.textbox(slide, x + Emu(int(col_w / 2)) - Inches(1.0),
                       y + Inches(1.2), Inches(2.0), Inches(0.5))
        mp = U.add_para(mt, first=True, align=PP_ALIGN.CENTER)
        mr = mp.add_run()
        U.style_run(mr, text=middle["value"], size=spec.get("middle_size", 20), bold=True,
                    color=U.color_hex(middle.get("color", theme["colors"]["accent"]), theme),
                    font=_font(spec, ctx, "heading"))


# ================================================================== 21. footer
def render_footer(slide, spec, theme, ctx):
    y = ctx["slide_height"] - Inches(0.45)
    text = spec.get("text", "")
    tf = U.textbox(slide, Inches(0.6), y, Inches(8), Inches(0.35))
    p = U.add_para(tf, first=True)
    r = p.add_run()
    U.style_run(r, text=text, size=10,
                color=_f(spec, "color", "text_muted", theme, "text_muted"),
                font=_font(spec, ctx))
    if spec.get("page", False) and spec.get("page_number"):
        pt = U.textbox(slide, ctx["slide_width"] - Inches(1.2), y, Inches(0.8), Inches(0.35))
        pp = U.add_para(pt, first=True, align=PP_ALIGN.RIGHT)
        pr_ = pp.add_run()
        U.style_run(pr_, text=str(spec["page_number"]), size=11, bold=True,
                    color=_f(spec, "page_color", "primary", theme, "primary"),
                    font=_font(spec, ctx))


# ================================================================== 22. page_number
def render_page_number(slide, spec, theme, ctx):
    y = ctx["slide_height"] - Inches(0.45)
    tf = U.textbox(slide, ctx["slide_width"] - Inches(1.2), y, Inches(0.8), Inches(0.35))
    p = U.add_para(tf, first=True, align=PP_ALIGN.RIGHT)
    r = p.add_run()
    U.style_run(r, text=str(spec.get("number", 1)), size=11, bold=True,
                color=_f(spec, "color", "primary", theme, "primary"),
                font=_font(spec, ctx))


# ================================================================== 线段绘制辅助
def _draw_segment(slide, x1, y1, x2, y2, color, alpha=90, thickness=0.03):
    """用旋转矩形绘制任意方向线段（中点+旋转法）。

    相比"起点+旋转"法，此方法保证：
      - 矩形宽高恒为正（避免负尺寸导致 PowerPoint 报"内容有问题"）
      - 旋转绕线段中点，位置准确不偏移
    """
    import math
    x1, y1, x2, y2 = int(x1), int(y1), int(x2), int(y2)
    dx, dy = x2 - x1, y2 - y1
    length = int(math.hypot(dx, dy))
    if length < 9144:
        return None
    midx, midy = (x1 + x2) // 2, (y1 + y2) // 2
    seg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(midx - length // 2), Emu(midy - int(914400 * thickness / 2)),
        Emu(length), Inches(thickness))
    U.set_shape_fill(seg, color, alpha)
    seg.line.fill.background()
    U.no_shadow(seg)
    seg.rotation = math.degrees(math.atan2(dy, dx))
    return seg


# ================================================================== 23. column_chart 竖向柱状图
def render_column_chart(slide, spec, theme, ctx):
    x, y, w, h = _geo(spec, ctx)
    data = spec.get("data", [])
    if not data:
        return
    color = _f(spec, "color", "primary", theme, "primary")
    color2 = _f(spec, "color2", "secondary", theme, "secondary")
    axis_color = _f(spec, "axis_color", "card_border", theme, "card_border")
    n = len(data)
    max_val = max([float(d.get("value", 0)) for d in data] + [1.0])
    gap = Inches(spec.get("gap", 0.45))
    col_w = Emu(max(int((w - gap * (n - 1)) / n), Inches(0.15)))
    chart_h = h - Inches(0.9)  # 顶部留数值，底部留标签
    base_y = y + chart_h
    # 基准线
    axis = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, base_y - Inches(0.02), w, Inches(0.03))
    U.set_shape_fill(axis, axis_color, 50)
    axis.line.fill.background()
    U.no_shadow(axis)
    for i, d in enumerate(data):
        cx = x + Emu(int(i * (col_w + gap)))
        frac = float(d.get("value", 0)) / max_val
        ch = Emu(max(int(chart_h * frac), Inches(0.05)))
        col = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, cx, base_y - ch, col_w, max(ch, Emu(9144)))
        if spec.get("gradient"):
            U.make_gradient_fill(col, [color, color2], angle=90)
        else:
            U.set_shape_fill(col, color if i % 2 == 0 else color2, spec.get("alpha"))
        U.set_round_rect_radius(col, 0.35)
        col.line.fill.background()
        U.no_shadow(col)
        # 数值（柱顶）
        vtf = U.textbox(slide, cx - Inches(0.2), base_y - ch - Inches(0.35), col_w + Inches(0.4), Inches(0.3))
        vp = U.add_para(vtf, first=True, align=PP_ALIGN.CENTER)
        vr = vp.add_run()
        U.style_run(vr, text=str(d.get("value", "")), size=12, bold=True,
                    color=_f(spec, "value_color", "text", theme, "text"),
                    font=_font(spec, ctx))
        # 标签（柱底）
        ltf = U.textbox(slide, cx - Inches(0.2), base_y + Inches(0.08), col_w + Inches(0.4), Inches(0.35))
        lp = U.add_para(ltf, first=True, align=PP_ALIGN.CENTER)
        lr = lp.add_run()
        U.style_run(lr, text=str(d.get("label", "")), size=12,
                    color=_f(spec, "label_color", "text_muted", theme, "text_muted"),
                    font=_font(spec, ctx))


# ================================================================== 24. line_chart 折线图
def render_line_chart(slide, spec, theme, ctx):
    import math
    x, y, w, h = _geo(spec, ctx)
    series = spec.get("series", [])
    if not series:
        return
    axis_color = _f(spec, "axis_color", "card_border", theme, "card_border")
    grid_color = _f(spec, "grid_color", "card_border", theme, "card_border")
    # 图表区
    pad_l = Inches(0.6)
    pad_r = Inches(0.3)
    pad_t = Inches(0.4)
    pad_b = Inches(0.5)
    chart_w = w - pad_l - pad_r
    chart_h = h - pad_t - pad_b
    plot_x = x + pad_l
    plot_y = y + pad_t
    # 收集所有点
    all_vals = []
    for s in series:
        for p in s.get("points", []):
            all_vals.append(float(p.get("value", 0)))
    max_val = max(all_vals + [1.0])
    min_val = min(all_vals + [0.0])
    span = max(max_val - min_val, 0.0001)
    n_points = max(len(series[0].get("points", [])) - 1, 1)
    # 网格线（4 条水平）
    for g in range(5):
        gy = plot_y + Emu(int(chart_h * g / 4))
        gl = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, plot_x, gy, chart_w, Inches(0.012))
        U.set_shape_fill(gl, grid_color, 25)
        gl.line.fill.background()
        U.no_shadow(gl)
        # y 轴刻度
        val = max_val - span * g / 4
        gtf = U.textbox(slide, x, gy - Inches(0.12), pad_l - Inches(0.1), Inches(0.25))
        gp = U.add_para(gtf, first=True, align=PP_ALIGN.RIGHT)
        gr = gp.add_run()
        U.style_run(gr, text="%.0f" % val, size=10,
                    color=_f(spec, "tick_color", "text_muted", theme, "text_muted"),
                    font=_font(spec, ctx))
    # 坐标轴
    ax_y = plot_y + chart_h
    ax = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, plot_x, ax_y - Inches(0.015), chart_w, Inches(0.025))
    U.set_shape_fill(ax, axis_color, 60)
    ax.line.fill.background()
    U.no_shadow(ax)
    # 绘制每条线
    palette = ["primary", "secondary", "accent", "success", "warning", "danger"]
    for si, s in enumerate(series):
        points = s.get("points", [])
        if len(points) < 2:
            continue
        scolor = U.color_hex(s.get("color", palette[si % len(palette)]), theme)
        pts = []
        for pi, p in enumerate(points):
            px = plot_x + Emu(int(chart_w * pi / n_points))
            py = plot_y + Emu(int(chart_h * (1 - (float(p.get("value", 0)) - min_val) / span)))
            pts.append((px, py))
            # 数据点
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, px - Inches(0.06), py - Inches(0.06),
                                         Inches(0.12), Inches(0.12))
            U.set_shape_fill(dot, scolor)
            dot.line.color.rgb = U.hex_to_rgb(_f(spec, "bg", "card_bg", theme, "card_bg"))
            dot.line.width = Pt(1.5)
            U.no_shadow(dot)
            # 数据标签
            if spec.get("labels", True) and pi % 2 == 0:
                ltf = U.textbox(slide, px - Inches(0.35), py - Inches(0.35), Inches(0.7), Inches(0.25))
                lp = U.add_para(ltf, first=True, align=PP_ALIGN.CENTER)
                lr = lp.add_run()
                U.style_run(lr, text=str(p.get("value", "")), size=10, bold=True,
                            color=scolor, font=_font(spec, ctx))
        # 连接线：用细长矩形旋转（中点法，避免负尺寸与偏移）
        for i in range(len(pts) - 1):
            x1, y1 = pts[i]
            x2, y2 = pts[i + 1]
            _draw_segment(slide, x1, y1, x2, y2, scolor, spec.get("line_alpha", 90), 0.035)
        # 图例
        if spec.get("legend", True) and s.get("label"):
            ly = y + Inches(0.05)
            ltf = U.textbox(slide, plot_x + chart_w - Inches(2.6), ly, Inches(2.6), Inches(0.3))
            lp = U.add_para(ltf, first=True, align=PP_ALIGN.RIGHT)
            lr = lp.add_run()
            U.style_run(lr, text="— " + s["label"], size=11, bold=True,
                        color=scolor, font=_font(spec, ctx))
    # X 轴标签
    labels = spec.get("labels") or [str(p.get("label", "")) for p in series[0].get("points", [])]
    for pi, lb in enumerate(labels):
        px = plot_x + Emu(int(chart_w * pi / n_points))
        ltf = U.textbox(slide, px - Inches(0.6), ax_y + Inches(0.05), Inches(1.2), Inches(0.3))
        lp = U.add_para(ltf, first=True, align=PP_ALIGN.CENTER)
        lr = lp.add_run()
        U.style_run(lr, text=str(lb), size=10,
                    color=_f(spec, "label_color", "text_muted", theme, "text_muted"),
                    font=_font(spec, ctx))


# ================================================================== 25. radar_chart 雷达图
def render_radar_chart(slide, spec, theme, ctx):
    import math
    x, y = U.parse_len(spec.get("x", 0), ctx["slide_width"]), U.parse_len(spec.get("y", 0), ctx["slide_height"])
    size = Inches(spec.get("size", 4.0))
    cx = x + Emu(int(size / 2))
    cy = y + Emu(int(size / 2))
    radius = Emu(int(size / 2) - Inches(0.2))
    labels = spec.get("labels", [])
    values = spec.get("values", [])
    n = len(labels)
    if n < 3:
        return
    grid_color = _f(spec, "grid_color", "card_border", theme, "card_border")
    # 网格：3 层多边形 + 连线
    for ring in range(1, 4):
        rr = Emu(int(radius * ring / 3))
        ring_pts = []
        for i in range(n):
            ang = -math.pi / 2 + 2 * math.pi * i / n
            ring_pts.append((cx + Emu(int(math.cos(ang) * rr)), cy + Emu(int(math.sin(ang) * rr))))
        # 多边形
        for i in range(n):
            x1, y1 = ring_pts[i]
            x2, y2 = ring_pts[(i + 1) % n]
            _draw_segment(slide, x1, y1, x2, y2, grid_color, 35, 0.02)
    # 轴
    for i in range(n):
        ang = -math.pi / 2 + 2 * math.pi * i / n
        ex = cx + Emu(int(math.cos(ang) * radius))
        ey = cy + Emu(int(math.sin(ang) * radius))
        _draw_segment(slide, cx, cy, ex, ey, grid_color, 45, 0.022)
        # 轴标签
        lx = cx + Emu(int(math.cos(ang) * (radius + Inches(0.35))))
        ly = cy + Emu(int(math.sin(ang) * (radius + Inches(0.35))))
        ltf = U.textbox(slide, lx - Inches(0.7), ly - Inches(0.15), Inches(1.4), Inches(0.3))
        lp = U.add_para(ltf, first=True, align=PP_ALIGN.CENTER)
        lr = lp.add_run()
        U.style_run(lr, text=labels[i], size=11, bold=True,
                    color=_f(spec, "label_color", "text", theme, "text"),
                    font=_font(spec, ctx))
    # 数据多边形（支持多个系列）
    palette = ["primary", "secondary", "accent", "success"]
    if isinstance(values[0], (list, tuple)):
        series = values
        labels_s = labels
    else:
        series = [values]
        labels_s = labels
    max_v = max([max(s if isinstance(s, list) else [s]) for s in series] + [1.0])
    for si, sv in enumerate(series):
        scolor = U.color_hex(spec.get("color", palette[si % len(palette)]), theme)
        if si > 0:
            scolor = U.color_hex(spec.get("color2", palette[si % len(palette)]), theme)
        pts = []
        for i, v in enumerate(sv):
            ang = -math.pi / 2 + 2 * math.pi * i / n
            rr = Emu(int(radius * float(v) / max_v))
            pts.append((cx + Emu(int(math.cos(ang) * rr)), cy + Emu(int(math.sin(ang) * rr))))
        # 填充多边形
        for i in range(n):
            x1, y1 = pts[i]
            x2, y2 = pts[(i + 1) % n]
            _draw_segment(slide, x1, y1, x2, y2, scolor, 90, 0.03)
        # 顶点
        for px, py in pts:
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, px - Inches(0.05), py - Inches(0.05),
                                         Inches(0.1), Inches(0.1))
            U.set_shape_fill(dot, scolor)
            dot.line.fill.background()
            U.no_shadow(dot)


# ================================================================== 26. funnel_chart 漏斗图
def render_funnel_chart(slide, spec, theme, ctx):
    x, y, w, h = _geo(spec, ctx)
    data = spec.get("data", [])
    if not data:
        return
    palette = ["primary", "secondary", "accent", "success", "warning", "danger"]
    max_val = max([float(d.get("value", 0)) for d in data] + [1.0])
    n = len(data)
    slot_h = Emu(int(h / n))
    gap = Inches(spec.get("gap", 0.12))
    for i, d in enumerate(data):
        frac = float(d.get("value", 0)) / max_val
        col_w = Emu(max(int(w * (0.35 + 0.65 * frac)), Inches(0.2)))
        col_x = x + Emu(int((w - col_w) / 2))
        col_y = y + Emu(int(i * slot_h)) + Emu(int(gap / 2))
        ch = Emu(max(int(slot_h - gap), Inches(0.2)))
        color = U.color_hex(d.get("color", palette[i % len(palette)]), theme)
        trap = slide.shapes.add_shape(MSO_SHAPE.TRAPEZOID, col_x, col_y, col_w, ch)
        U.set_shape_fill(trap, color, spec.get("alpha"))
        trap.line.fill.background()
        U.no_shadow(trap)
        # 文本（左侧标签 + 右侧值）
        tf = slide.shapes.add_textbox(col_x + Inches(0.25), col_y, col_w - Inches(0.5), ch)
        ttf = tf.text_frame
        ttf.word_wrap = True
        ttf.vertical_anchor = MSO_ANCHOR.MIDDLE
        ttf.margin_left = ttf.margin_right = 0
        ttf.margin_top = ttf.margin_bottom = 0
        tp = ttf.paragraphs[0]
        tr = tp.add_run()
        label = str(d.get("label", ""))
        value = str(d.get("value", ""))
        show = f"{label}  {value}" if value else label
        fg = "FFFFFF" if theme["style"] == "dark" else "FFFFFF"
        U.style_run(tr, text=show, size=14, bold=True, color=fg, font=_font(spec, ctx))
        # 转化率
        if i > 0 and spec.get("rate", True):
            prev = float(data[i - 1].get("value", 0))
            cur = float(d.get("value", 0))
            if prev > 0:
                rate = cur / prev * 100
                rtf = U.textbox(slide, col_x + col_w + Inches(0.15),
                                col_y + Emu(int(ch / 2)) - Inches(0.15), Inches(1.2), Inches(0.3))
                rp = U.add_para(rtf, first=True)
                rr = rp.add_run()
                U.style_run(rr, text=f"{rate:.0f}%", size=12, bold=True,
                            color=_f(spec, "rate_color", "accent", theme, "accent"),
                            font=_font(spec, ctx))


# ================================================================== 27. gauge 仪表盘
def render_gauge(slide, spec, theme, ctx):
    import math
    x, y = U.parse_len(spec.get("x", 0), ctx["slide_width"]), U.parse_len(spec.get("y", 0), ctx["slide_height"])
    size = Inches(spec.get("size", 3.0))
    value = max(0.0, min(1.0, float(spec.get("value", 0.7))))
    color = _f(spec, "color", "primary", theme, "primary")
    track = _f(spec, "track", "card_border", theme, "card_border")
    cx = x + Emu(int(size / 2))
    cy = y + Emu(int(size * 0.85))
    radius = Emu(int(size * 0.42))
    # 轨道（半圆扇形，用 5 段近似）
    n_seg = 10
    sweep = math.pi
    start_ang = math.pi
    # 轨道
    for i in range(n_seg):
        a1 = start_ang + sweep * i / n_seg
        a2 = start_ang + sweep * (i + 1) / n_seg
        frac_done = i / n_seg
        is_filled = frac_done <= value
        c = color if is_filled else track
        if not is_filled:
            seg = slide.shapes.add_shape(
                MSO_SHAPE.PIE,
                cx - radius, cy - radius, Emu(int(radius * 2)), Emu(int(radius * 2)))
            seg.adjustments[0] = a1 / (2 * math.pi)
            seg.adjustments[1] = a2 / (2 * math.pi)
            U.set_shape_fill(seg, c, 40)
            seg.line.color.rgb = U.hex_to_rgb(_f(spec, "bg", "card_bg", theme, "card_bg"))
            seg.line.width = Pt(1)
            U.no_shadow(seg)
        else:
            seg = slide.shapes.add_shape(
                MSO_SHAPE.PIE,
                cx - radius, cy - radius, Emu(int(radius * 2)), Emu(int(radius * 2)))
            seg.adjustments[0] = a1 / (2 * math.pi)
            seg.adjustments[1] = a2 / (2 * math.pi)
            if spec.get("gradient"):
                U.make_gradient_fill(seg, [color, _f(spec, "color2", "secondary", theme, "secondary")], 90)
            else:
                U.set_shape_fill(seg, color)
            seg.line.color.rgb = U.hex_to_rgb(_f(spec, "bg", "card_bg", theme, "card_bg"))
            seg.line.width = Pt(1)
            U.no_shadow(seg)
    # 中心数值
    vtf = U.textbox(slide, cx - Inches(0.8), cy - Inches(0.55), Inches(1.6), Inches(0.6),
                    anchor=MSO_ANCHOR.MIDDLE)
    vp = U.add_para(vtf, first=True, align=PP_ALIGN.CENTER)
    vr = vp.add_run()
    U.style_run(vr, text=str(spec.get("value_text", f"{int(value * 100)}%")),
                size=spec.get("value_size", 30), bold=True, color=color,
                font=_font(spec, ctx, "heading"))
    if spec.get("label"):
        ltf = U.textbox(slide, cx - Inches(1.2), cy + Inches(0.05), Inches(2.4), Inches(0.35))
        lp = U.add_para(ltf, first=True, align=PP_ALIGN.CENTER)
        lr = lp.add_run()
        U.style_run(lr, text=spec["label"], size=12,
                    color=_f(spec, "label_color", "text_muted", theme, "text_muted"),
                    font=_font(spec, ctx))


# ================================================================== 28. pricing 定价卡
def render_pricing(slide, spec, theme, ctx):
    x, y, w, h = _geo(spec, ctx)
    plans = spec.get("plans", [])
    if not plans:
        return
    n = len(plans)
    gap = Inches(spec.get("gap", 0.3))
    card_w = Emu(int((w - gap * (n - 1)) / n))
    highlight_idx = spec.get("highlight", -1)
    for i, plan in enumerate(plans):
        cx = x + Emu(int(i * (card_w + gap)))
        highlight = (i == highlight_idx)
        bg = _f(spec, "bg", "card_bg", theme, "card_bg")
        accent = U.color_hex(plan.get("color", theme["colors"]["primary"]), theme)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, y, card_w, h)
        if highlight:
            U.make_gradient_fill(card, [accent, _f(spec, "color2", "secondary", theme, "secondary")], 90)
            U.add_shadow(card, blur=0.16, dist=0.05, alpha=40)
        else:
            U.set_shape_fill(card, bg)
            U.set_round_rect_radius(card, spec.get("radius", 0.08))
            U.set_shape_line(card, _f(spec, "border", "card_border", theme, "card_border"), 1)
            U.no_shadow(card)
        U.set_round_rect_radius(card, spec.get("radius", 0.08))
        pad = Inches(0.35)
        tf = U.textbox(slide, cx + pad, y + pad, card_w - 2 * pad, h - 2 * pad)
        # 热门标签
        if plan.get("tag"):
            tag = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, cx + card_w - Inches(1.6), y + Inches(0.15),
                Inches(1.3), Inches(0.35))
            U.set_shape_fill(tag, accent, 90)
            tag.line.fill.background()
            U.no_shadow(tag)
            ttf = tag.text_frame
            ttf.word_wrap = False
            ttf.vertical_anchor = MSO_ANCHOR.MIDDLE
            ttf.margin_left = ttf.margin_right = Inches(0.05)
            ttf.margin_top = ttf.margin_bottom = 0
            tp = ttf.paragraphs[0]
            tp.alignment = PP_ALIGN.CENTER
            tr = tp.add_run()
            U.style_run(tr, text=plan["tag"], size=10, bold=True,
                        color="FFFFFF", font=_font(spec, ctx))
        # 名称
        p0 = U.add_para(tf, first=True)
        r0 = p0.add_run()
        fg = "FFFFFF" if (highlight or theme["style"] == "dark") else theme["colors"]["text"]
        U.style_run(r0, text=plan.get("name", ""), size=18, bold=True, color=fg,
                    font=_font(spec, ctx, "heading"))
        # 价格
        p1 = U.add_para(tf, space_before=8)
        r1 = p1.add_run()
        U.style_run(r1, text=plan.get("price", ""), size=34, bold=True,
                    color="FFFFFF" if highlight else accent, font=_font(spec, ctx, "heading"))
        if plan.get("per"):
            r1b = p1.add_run()
            U.style_run(r1b, text=" / " + plan["per"], size=13,
                        color="FFFFFF" if highlight else _f(spec, "muted_color", "text_muted", theme, "text_muted"),
                        font=_font(spec, ctx))
        # 特性列表
        for feat in plan.get("features", []):
            fp = U.add_para(tf, space_before=4, line_spacing=1.25)
            fr = fp.add_run()
            U.style_run(fr, text="✓ " + feat, size=12,
                        color="FFFFFF" if highlight else _f(spec, "feature_color", "text", theme, "text"),
                        font=_font(spec, ctx))
        # 按钮
        if plan.get("button"):
            btn_h = Inches(0.45)
            btn = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                cx + pad, y + h - btn_h - Inches(0.3), card_w - 2 * pad, btn_h)
            if highlight:
                U.set_shape_fill(btn, "FFFFFF")
                U.style_run(btn.text_frame.paragraphs[0].add_run(), text=plan["button"],
                            size=13, bold=True, color=accent, font=_font(spec, ctx))
            else:
                U.set_shape_fill(btn, accent, 20)
                U.style_run(btn.text_frame.paragraphs[0].add_run(), text=plan["button"],
                            size=13, bold=True, color=accent, font=_font(spec, ctx))
            U.set_round_rect_radius(btn, 0.5)
            btn.line.fill.background()
            U.no_shadow(btn)
            btf = btn.text_frame
            btf.vertical_anchor = MSO_ANCHOR.MIDDLE
            btf.margin_left = btf.margin_right = 0
            btf.margin_top = btf.margin_bottom = 0
            bp = btf.paragraphs[0]
            bp.alignment = PP_ALIGN.CENTER


# ================================================================== 29. flow 流程图
def render_flow(slide, spec, theme, ctx):
    x, y, w, h = _geo(spec, ctx)
    nodes = spec.get("nodes", [])
    if not nodes:
        return
    direction = spec.get("direction", "horizontal")  # horizontal | vertical
    gap = Inches(spec.get("gap", 0.5))
    n = len(nodes)
    if direction == "vertical":
        slot_h = Emu(int((h - gap * (n - 1)) / n))
        for i, nd in enumerate(nodes):
            ny = y + Emu(int(i * (slot_h + gap)))
            render_shape(slide, {
                "shape": nd.get("shape", "round_rect"), "x": x, "y": ny, "w": w, "h": slot_h,
                "fill": nd.get("fill", "gradient"), "fill_color": nd.get("color"),
                "text": nd.get("text"), "text_size": nd.get("text_size", 16),
                "text_color": "text_invert", "radius": nd.get("radius", 0.12),
            }, theme, ctx)
            if i < n - 1:
                ar = slide.shapes.add_shape(
                    MSO_SHAPE.DOWN_ARROW,
                    x + Emu(int((w - Inches(0.3)) / 2)),
                    ny + slot_h + Inches(0.02), Inches(0.3), Inches(0.35))
                U.set_shape_fill(ar, _f(spec, "arrow_color", "primary", theme, "primary"), 70)
                ar.line.fill.background()
                U.no_shadow(ar)
    else:
        slot_w = Emu(int((w - gap * (n - 1)) / n))
        for i, nd in enumerate(nodes):
            nx = x + Emu(int(i * (slot_w + gap)))
            render_shape(slide, {
                "shape": nd.get("shape", "round_rect"), "x": nx, "y": y, "w": slot_w, "h": h,
                "fill": nd.get("fill", "gradient"), "fill_color": nd.get("color"),
                "text": nd.get("text"), "text_size": nd.get("text_size", 16),
                "text_color": "text_invert", "radius": nd.get("radius", 0.12),
            }, theme, ctx)
            if i < n - 1:
                ar = slide.shapes.add_shape(
                    MSO_SHAPE.RIGHT_ARROW,
                    nx + slot_w + Inches(0.02), y + Emu(int((h - Inches(0.3)) / 2)),
                    Inches(0.4), Inches(0.3))
                U.set_shape_fill(ar, _f(spec, "arrow_color", "primary", theme, "primary"), 70)
                ar.line.fill.background()
                U.no_shadow(ar)


# ================================================================== 30. cloud 标签云
def render_cloud(slide, spec, theme, ctx):
    x, y, w, h = _geo(spec, ctx)
    tags = spec.get("tags", [])
    if not tags:
        return
    # 简单流式布局
    cur_x = x
    cur_y = y
    max_y = y + h
    row_h = Inches(0.5)
    for tag in tags:
        if isinstance(tag, str):
            tag = {"text": tag}
        text = tag.get("text", "")
        size = tag.get("size", 13)
        color = U.color_hex(tag.get("color", "primary"), theme)
        weight = tag.get("weight", 1)
        pad_x = Inches(0.25)
        est_w = max(Inches(0.8), Inches(len(text) * size * 0.17 + 0.4))
        # 换行
        if cur_x + est_w > x + w:
            cur_x = x
            cur_y += row_h
        if cur_y + row_h > max_y:
            break
        tag_sh = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, cur_x, cur_y, est_w, Inches(0.42))
        U.set_round_rect_radius(tag_sh, 0.5)
        if tag.get("variant", "soft") == "soft":
            U.set_shape_fill(tag_sh, color, 16 + 6 * weight)
            tag_sh.line.fill.background()
            txt_color = color
        elif tag.get("variant") == "fill":
            U.set_shape_fill(tag_sh, color)
            tag_sh.line.fill.background()
            txt_color = "FFFFFF"
        else:
            U.set_shape_fill(tag_sh, None)
            U.set_shape_line(tag_sh, color, 1.2)
            txt_color = color
        U.no_shadow(tag_sh)
        ttf = tag_sh.text_frame
        ttf.word_wrap = False
        ttf.vertical_anchor = MSO_ANCHOR.MIDDLE
        ttf.margin_left = ttf.margin_right = Inches(0.1)
        ttf.margin_top = ttf.margin_bottom = 0
        tp = ttf.paragraphs[0]
        tp.alignment = PP_ALIGN.CENTER
        tr = tp.add_run()
        U.style_run(tr, text=text, size=size, bold=(weight >= 2), color=txt_color,
                    font=_font(spec, ctx))
        cur_x += est_w + Inches(0.15)


# ================================================================== 31. carousel 图片轮播/画廊
def render_carousel(slide, spec, theme, ctx):
    x, y, w, h = _geo(spec, ctx)
    images = spec.get("images", [])
    if not images:
        return
    cols = int(spec.get("cols", 3))
    gap = Inches(spec.get("gap", 0.2))
    img_w = Emu(int((w - gap * (cols - 1)) / cols))
    rows = (len(images) + cols - 1) // cols
    img_h = spec.get("img_height", 2.2)
    if isinstance(img_h, (int, float)) and img_h <= 10:
        img_h = Inches(img_h)
    for i, img in enumerate(images):
        r_i, c_i = divmod(i, cols)
        cx = x + Emu(int(c_i * (img_w + gap)))
        cy = y + Emu(int(r_i * (img_h + gap)))
        path = img.get("path")
        # 占位色块（无图时）
        ph = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, cy, img_w, img_h)
        U.set_shape_fill(ph, U.color_hex(img.get("color", "card_bg2"), theme))
        U.set_round_rect_radius(ph, spec.get("radius", 0.06))
        U.no_shadow(ph)
        if path and os.path.isfile(path):
            try:
                pic = slide.shapes.add_picture(path, cx, cy, width=img_w)
                if pic.height > img_h:
                    pic.height = img_h
                if spec.get("radius"):
                    U.set_rounded_picture(pic, spec["radius"])
            except Exception:
                pass
        else:
            # 占位图标
            icon = img.get("icon", "🖼️")
            phtf = ph.text_frame
            phtf.word_wrap = True
            phtf.vertical_anchor = MSO_ANCHOR.MIDDLE
            phtf.margin_left = phtf.margin_right = 0
            phtf.margin_top = phtf.margin_bottom = 0
            php = phtf.paragraphs[0]
            php.alignment = PP_ALIGN.CENTER
            phr = php.add_run()
            U.style_run(phr, text=icon, size=30, font=_font(spec, ctx))
        # 底部标题条
        if img.get("caption"):
            cap_h = Inches(0.4)
            cap = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, cx, cy + img_h - cap_h, img_w, cap_h)
            U.set_shape_fill(cap, "000000", 55)
            cap.line.fill.background()
            U.no_shadow(cap)
            ctf = cap.text_frame
            ctf.word_wrap = True
            ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
            ctf.margin_left = Inches(0.12)
            ctf.margin_right = Inches(0.05)
            ctf.margin_top = ctf.margin_bottom = 0
            cp = ctf.paragraphs[0]
            cr = cp.add_run()
            U.style_run(cr, text=img["caption"], size=11, bold=True,
                        color="FFFFFF", font=_font(spec, ctx))


# ================================================================== 32. highlight 高亮数字
def render_highlight(slide, spec, theme, ctx):
    x, y, w, h = _geo(spec, ctx)
    items = spec.get("items", [])
    if not items:
        return
    n = len(items)
    gap = Inches(spec.get("gap", 0.4))
    slot_w = Emu(int((w - gap * (n - 1)) / n))
    for i, item in enumerate(items):
        cx = x + Emu(int(i * (slot_w + gap)))
        color = U.color_hex(item.get("color", "primary"), theme)
        # 大数字
        vtf = U.textbox(slide, cx, y, slot_w, Inches(1.2))
        vp = U.add_para(vtf, first=True, align=PP_ALIGN.CENTER)
        vr = vp.add_run()
        U.style_run(vr, text=str(item.get("value", "")), size=item.get("size", 40),
                    bold=True, color=color, font=_font(spec, ctx, "heading"))
        if item.get("unit"):
            ur = vp.add_run()
            U.style_run(ur, text=item["unit"], size=item.get("unit_size", 16),
                        color=_f(spec, "unit_color", "text_muted", theme, "text_muted"),
                        font=_font(spec, ctx))
        # 标签
        ltf = U.textbox(slide, cx, y + Inches(1.1), slot_w, Inches(0.4))
        lp = U.add_para(ltf, first=True, align=PP_ALIGN.CENTER)
        lr = lp.add_run()
        U.style_run(lr, text=str(item.get("label", "")), size=13,
                    color=_f(spec, "label_color", "text_muted", theme, "text_muted"),
                    font=_font(spec, ctx))
        # 底部强调线
        line = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            cx + Emu(int((slot_w - Inches(0.5)) / 2)), y + Inches(1.6),
            Inches(0.5), Inches(0.06))
        U.set_shape_fill(line, color)
        U.set_round_rect_radius(line, 0.5)
        line.line.fill.background()
        U.no_shadow(line)


# ================================================================== 注册表
COMPONENTS = {
    "title": render_title,
    "text": render_text,
    "list": render_list,
    "card": render_card,
    "kpi": render_kpi,
    "badge": render_badge,
    "divider": render_divider,
    "shape": render_shape,
    "icon": render_icon,
    "progress": render_progress,
    "bar_chart": render_bar_chart,
    "column_chart": render_column_chart,
    "line_chart": render_line_chart,
    "radar_chart": render_radar_chart,
    "funnel_chart": render_funnel_chart,
    "gauge": render_gauge,
    "pie_chart": render_pie_chart,
    "pricing": render_pricing,
    "flow": render_flow,
    "cloud": render_cloud,
    "carousel": render_carousel,
    "highlight": render_highlight,
    "timeline": render_timeline,
    "steps": render_steps,
    "grid_cards": render_grid_cards,
    "table": render_table,
    "quote": render_quote,
    "image": render_image,
    "team": render_team,
    "stat_compare": render_stat_compare,
    "footer": render_footer,
    "page_number": render_page_number,
}


def render_component(slide, spec, theme, ctx):
    """渲染单个组件。未知类型静默跳过。"""
    ctype = str(spec.get("type", "")).lower()
    fn = COMPONENTS.get(ctype)
    if fn is None:
        return False
    try:
        fn(slide, spec, theme, ctx)
        return True
    except Exception:
        # 单组件失败不影响整页
        import traceback
        traceback.print_exc()
        return False


def component_list():
    """返回组件清单（供 /api/components）。"""
    return [
        {"type": "title", "desc": "标题（可带强调条）", "fields": ["text", "size", "color", "align", "accent_bar"]},
        {"type": "text", "desc": "富文本段落（多段多 run）", "fields": ["content/paragraphs", "size", "color", "align", "line_spacing"]},
        {"type": "list", "desc": "列表（bullet/number/icon）", "fields": ["items", "marker", "marker_color"]},
        {"type": "card", "desc": "卡片（图标+标题+文本+阴影）", "fields": ["icon", "title", "text", "accent", "bg", "radius"]},
        {"type": "kpi", "desc": "KPI 指标卡（数值+单位+增量）", "fields": ["value", "unit", "label", "delta", "icon", "accent"]},
        {"type": "badge", "desc": "标签（fill/outline/soft）", "fields": ["text", "variant", "color"]},
        {"type": "divider", "desc": "分隔线（line/gradient/dot）", "fields": ["style", "color"]},
        {"type": "shape", "desc": "自定义形状（圆角/箭头/多边形等）", "fields": ["shape", "fill", "line", "text", "rotation"]},
        {"type": "icon", "desc": "Emoji 图标（圆形底）", "fields": ["emoji", "size", "color"]},
        {"type": "progress", "desc": "进度条", "fields": ["value", "color", "label", "gradient"]},
        {"type": "bar_chart", "desc": "横向条形图", "fields": ["data[{label,value}]", "color", "gradient"]},
        {"type": "column_chart", "desc": "竖向柱状图（渐变）", "fields": ["data[{label,value}]", "color", "gradient"]},
        {"type": "line_chart", "desc": "折线图（多系列/数据点/网格）", "fields": ["series[{label,points}]", "labels"]},
        {"type": "radar_chart", "desc": "雷达图（多系列）", "fields": ["labels", "values", "size"]},
        {"type": "funnel_chart", "desc": "漏斗图（含转化率）", "fields": ["data[{label,value,color}]", "rate"]},
        {"type": "gauge", "desc": "仪表盘（半圆进度）", "fields": ["value", "value_text", "label", "size"]},
        {"type": "pie_chart", "desc": "饼图/环形图 + 图例", "fields": ["data[{label,value,color}]", "donut", "center_text"]},
        {"type": "pricing", "desc": "定价卡（多方案/高亮/按钮）", "fields": ["plans[{name,price,features,button}]", "highlight"]},
        {"type": "flow", "desc": "流程图（横向/纵向 + 箭头）", "fields": ["nodes[{text,shape,color}]", "direction"]},
        {"type": "cloud", "desc": "标签云（软/实心/描边）", "fields": ["tags[{text,size,color,weight}]"]},
        {"type": "carousel", "desc": "图片画廊（圆角/标题条）", "fields": ["images[{path,caption,color}]", "cols"]},
        {"type": "highlight", "desc": "高亮数字组", "fields": ["items[{value,unit,label,color}]"]},
        {"type": "timeline", "desc": "时间线（纵向/横向）", "fields": ["items[{date,title,desc}]", "vertical"]},
        {"type": "steps", "desc": "步骤流程条", "fields": ["items", "current"]},
        {"type": "grid_cards", "desc": "卡片网格", "fields": ["cards", "cols", "gap"]},
        {"type": "table", "desc": "增强表格（表头配色/斑马纹/首列强调）", "fields": ["headers", "rows", "header_bg"]},
        {"type": "quote", "desc": "引用块", "fields": ["text", "author", "accent"]},
        {"type": "image", "desc": "图片（可圆角/描边/阴影）", "fields": ["path", "radius", "border"]},
        {"type": "team", "desc": "团队成员卡", "fields": ["members[{name,role,avatar,color}]"]},
        {"type": "stat_compare", "desc": "左右对比统计（VS）", "fields": ["left", "right", "middle"]},
        {"type": "footer", "desc": "页脚文字", "fields": ["text", "page"]},
        {"type": "page_number", "desc": "页码", "fields": ["number"]},
    ]
